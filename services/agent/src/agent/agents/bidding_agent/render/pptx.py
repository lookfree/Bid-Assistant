from __future__ import annotations
import io
import logging
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.enum.shapes import PP_PLACEHOLDER
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt
from agent.agents.bidding_agent.render.pptx_blocks import (
    SLIDE_H, SLIDE_W, accent_bar, body_band, bullets_box, content_box,
    paint_surface, rect, render_cover, rounded_rect, slide_size, textbox, title_row,
)
from agent.agents.bidding_agent.render.styles import (
    CARD, LEAD, TYPE, blend_toward, chart_palette, is_dark, on_primary, tokens_for,
)
from agent.agents.bidding_agent.schemas import DeckSpec, Slide, SlideChart, StatItem

logger = logging.getLogger(__name__)

# 本模块只管 deck 的组装与「每套模板长得一样」的那些部件（图表/数字卡/评分点角标/页码/
# 分隔页/结束页）。模板之间真正不同的封面、标题行、要点卡三族画法在 render/pptx_blocks.py，
# 配色与结构开关在 render/styles.py。企业自有母版走 render_pptx(master_bytes=...)：
# 强调色/评分点角标/页码仍取这套 token，让自绘部分和母版主题不违和。


def _chip_width(text: str) -> int:
    """评分点角标自适应宽度。
    原按 0.11in/字符估算，中文下少算三分之一：12pt 中文字符实际约 0.167in 宽，30 字标签
    估 3.3in、实需 ~4.7in——文字比框宽 1.4in，而形状文字默认居中且 word_wrap=False，
    多出的部分往两边各溢出 0.7in，左边正好顶着页边距把「评」字推出画面（用户实测每页都被裁）。
    改为按字符类型估宽（CJK/全角 0.17in、ASCII 0.085in）再留 0.4in 内边距。"""
    w = sum(0.17 if ord(c) > 0x2E80 else 0.085 for c in text)
    return int(Inches(max(2.5, min(9.0, w + 0.4))))


def _scoring_chip(slide, scoring: str, tokens: dict) -> None:
    """底部左侧圆角矩形评分点角标：浅底色 + 细强调边框 + 强调色文字。
    宽度夹到正文可用宽以内、纵向按页高定位——死写 6.55in 在页高不同的客户母版上会漂。"""
    text = f"评分点｜{scoring}"
    left, content_w, sh = content_box(slide)
    width = min(_chip_width(text), content_w)
    shape = rounded_rect(slide, left, sh - Inches(0.95), width, Inches(0.5), tokens["tint"],
                         line_rgb=tokens["accent"])
    tf = shape.text_frame
    tf.word_wrap = False
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Inches(0.15)
    tf.margin_right = Inches(0.15)
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = PP_ALIGN.LEFT   # 居中会让溢出往两边跑，左侧首字被推出画面
    run = p.runs[0] if p.runs else p.add_run()
    run.font.size = Pt(TYPE["note"])
    run.font.color.rgb = tokens["accent"]


def _page_number(slide, n: int, total: int, tokens: dict) -> None:
    """底部右侧页码 “n / total”，10pt 弱化灰。
    位置按幻灯片真实尺寸算：原来死写 16:9 常量，在 4:3 客户母版（10in 宽）上每一页的页码都被
    推出页面右侧近 3in（三种版式都带页码，问题从"偶发"变成"每页必现"）。"""
    sw, sh = slide_size(slide)
    textbox(slide, sw - Inches(1.9), sh - Inches(0.85), Inches(1.2), Inches(0.4),
            [f"{n} / {total}"], size=TYPE["caption"], color=tokens["muted"], align=PP_ALIGN.RIGHT)


def _render_section(slide, s: Slide, tokens: dict, n: int, total: int, seq: int = 1) -> None:
    """章节分隔页（述标结构性升级）：满屏主色块 + 居中大标题 + 可选一句过渡副标题（取 bullets[0]）。
    评审实测教训：所有正文页长得一模一样，评委翻到第 8 页都不知道"讲到哪个部分了"——
    按评分维度分组（项目理解/技术方案/团队业绩/服务承诺与报价/风险防控）时每组开头插一张，
    给整套述标制造视觉节奏。不对应具体评分点，不挂 scoring 角标。
    深色模板换一套底：它的 primary 是亮青绿，整页铺满就是一张刺眼的荧光页夹在深色正文页之间
    （实测投影下更糟），改用卡片底色当分隔页的地，强调色只留给短线和序号。"""
    sw, sh = slide_size(slide)
    dark = is_dark(tokens)
    ground = tokens["tint"] if dark else tokens["primary"]
    fg = tokens["text"] if dark else tokens["white"]
    rect(slide, 0, 0, sw, sh, ground)
    # 大号序号：只有一行居中标题时整页太空（用户反馈「太素」），用一个压在标题上方的
    # 淡色大数字撑住版面，同时给评委一个"第几部分"的强锚点。
    textbox(slide, Inches(1.2), sh * 0.16, Inches(4.0), Inches(2.1),
            [f"{seq:02d}"], size=TYPE["section_num"], color=blend_toward(ground, fg, 0.34),
            bold=True, line_spacing=LEAD["title"])
    rect(slide, Inches(1.2), sh * 0.545, Inches(0.9), Pt(3), tokens["accent"])
    textbox(slide, Inches(1.2), sh * 0.575, sw - Inches(2.4), Inches(1.2),
            [s.title], size=TYPE["section_title"], color=fg, bold=True,
            line_spacing=LEAD["title"])
    if s.bullets:
        # 同图表页：分隔页的过渡语也全部渲染，不静默丢弃第二句
        textbox(slide, Inches(1.2), sh * 0.735, sw - Inches(2.4), Inches(0.8),
                s.bullets, size=TYPE["section_lead"],
                color=blend_toward(fg, tokens["accent"], 0.25), line_spacing=LEAD["note"])
    _page_number(slide, n, total, tokens)


_CHART_TYPE_MAP = {
    "column": XL_CHART_TYPE.COLUMN_CLUSTERED,
    "bar": XL_CHART_TYPE.BAR_CLUSTERED,
    "pie": XL_CHART_TYPE.PIE,
    "line": XL_CHART_TYPE.LINE_MARKERS,
}


def _paint_pie(plot, tokens: dict, n_categories: int) -> None:
    """饼图逐扇区（point）上色；数值标签压在扇区上，整张图统一用 on_primary 当字色。
    为什么不逐扇区配字色、也不把标签移到扇区外：这两种写法 PowerPoint 认，预览用的渲染器
    都不认（实测逐扇区字色被整体忽略、outEnd 照样画在扇区里）。既然只能有一种字色，
    就把这一种字色交给 chart_palette——它逐块把扇区压到与标签色 ≥4.5:1，
    8 类之内既不重复取色，标签也一直读得出来（报价构成、岗位分布这类 7 类饼图很常见）。"""
    colors = chart_palette(tokens, n_categories, label_rgb=on_primary(tokens))
    for i, point in enumerate(plot.series[0].points):
        point.format.fill.solid()
        point.format.fill.fore_color.rgb = colors[i % len(colors)]
    plot.data_labels.font.color.rgb = on_primary(tokens)


def _chart_title(gchart, chart: SlideChart, tokens: dict) -> None:
    """单系列图（饼图/单柱图）显式写图表标题并配色。
    两个渲染器都会给单系列图自动补一个黑色标题，而"不要自动标题"的开关只有 PowerPoint 认，
    预览渲染器照画不误——深色模板下就是一行黑压黑的字。干脆自己写一个受控的。"""
    if len(chart.series) != 1:
        gchart.has_title = False
        return
    gchart.has_title = True
    p = gchart.chart_title.text_frame.paragraphs[0]
    p.text = chart.series[0].name
    run = p.runs[0] if p.runs else p.add_run()
    run.font.size = Pt(TYPE["chart_text"])
    run.font.bold = False
    run.font.color.rgb = tokens["muted"]


def _chart_legend(gchart, chart: SlideChart, tokens: dict) -> None:
    """图例位置与配色。饼图的图例放右侧：饼的直径受页高限制，图例摆下面等于把左右两侧
    白白空出来（用户实评「图表页左右偏空」）；摆右侧则由它把横向空间吃掉，整块图是满的。
    柱/线图本来就横向铺开，图例仍走下方。
    图例字色必须单独设：它不继承图表级 txPr，深色模板下整排图例是黑字（实测糊在底色里）。"""
    if not gchart.has_legend:
        return
    gchart.legend.position = (XL_LEGEND_POSITION.RIGHT if chart.type == "pie"
                              else XL_LEGEND_POSITION.BOTTOM)
    gchart.legend.include_in_layout = False
    gchart.legend.font.size = Pt(TYPE["chart_text"])
    gchart.legend.font.color.rgb = tokens["text"]


def _render_chart_body(slide, chart: SlideChart, tokens: dict, box: tuple) -> None:
    """图表页主体：真实 PowerPoint 图表对象（python-pptx add_chart），不是图片——评委可在
    PowerPoint 里直接编辑数值/改样式，这是"原生深度"而非"糊一张图上去"的关键差异。
    数据标签常开：评委扫一眼数字就懂，不用眯眼看坐标轴刻度。
    几何由调用方给（box=(left, top, width, height)）：有结论侧栏时图表只占主区，没有时占满。"""
    data = CategoryChartData()
    data.categories = chart.categories
    for series in chart.series:
        data.add_series(series.name, series.values)
    chart_type = _CHART_TYPE_MAP.get(chart.type, XL_CHART_TYPE.COLUMN_CLUSTERED)
    frame = slide.shapes.add_chart(chart_type, *(int(v) for v in box), data)
    gchart = frame.chart
    # 坐标轴文字统一走模板正文色：Office 默认黑字，深色模板下整套刻度直接黑压黑看不见
    gchart.font.color.rgb = tokens["text"]
    _chart_title(gchart, chart, tokens)
    plot = gchart.plots[0]
    plot.has_data_labels = True
    plot.data_labels.font.size = Pt(TYPE["chart_text"])
    plot.data_labels.font.color.rgb = tokens["text"]
    if chart.type == "pie":
        _paint_pie(plot, tokens, len(chart.categories))
        gchart.has_legend = True
    else:
        colors = chart_palette(tokens, len(chart.series))
        for i, series in enumerate(plot.series):
            series.format.fill.solid()
            series.format.fill.fore_color.rgb = colors[i % len(colors)]
        gchart.has_legend = len(chart.series) > 1
    _chart_legend(gchart, chart, tokens)


_RAIL_RATIO = 0.32       # 结论侧栏占正文可用宽的比例


def _chart_rail(slide, notes: list[str], tokens: dict, box: tuple) -> None:
    """图表页的「关键结论」侧栏：浅底面板 + 栏目名 + 细线 + 逐条结论。
    内容**只取 slide 已有的 bullets**（提示词里图表页本来就允许带 1-2 句结论），
    一个字都不额外编；没有结论就根本不画这一栏，让图表放大占满整幅。"""
    left, top, width, bottom = box
    pad = CARD["pad_x"]
    rounded_rect(slide, left, top, width, bottom - top, tokens["tint"])
    textbox(slide, left + pad, top + Inches(0.24), width - 2 * pad, Inches(0.3), ["关键结论"],
            size=TYPE["eyebrow"], color=tokens["accent"], bold=True, tracking=1.5)
    rule_y = top + Inches(0.62)
    rect(slide, left + pad, rule_y, width - 2 * pad, Pt(1),
         blend_toward(tokens["tint"], tokens["text"], 0.30))
    cells = _rail_cells(len(notes), rule_y + Inches(0.2), bottom - Inches(0.2))
    divider = blend_toward(tokens["tint"], tokens["text"], 0.18)
    for i, ((y, h), note) in enumerate(zip(cells, notes)):
        if i:
            rect(slide, left + pad, y, width - 2 * pad, Pt(0.75), divider)
        rect(slide, left + pad, y + h // 2 - Inches(0.06), Inches(0.12), Inches(0.12),
             tokens["accent"])
        tf = textbox(slide, left + pad + Inches(0.3), y, width - 2 * pad - Inches(0.3), h, [note],
                     size=TYPE["body"], color=tokens["text"], line_spacing=LEAD["body"])
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE


def _rail_cells(n: int, top: int, bottom: int) -> list[tuple[int, int]]:
    """侧栏 n 条结论的 (起始 y, 行高)：**均分整条可用高度**，每条在自己那格里垂直居中。
    只按行高排、贴着上沿走的话，1 条结论会缩在面板顶上、下面空掉四分之三块底板——
    那比不画侧栏还难看。均分之后条数再少也是"有意留白"，不是"没排满"。"""
    h = int((bottom - top) / max(1, n))
    return [(top + i * h, h) for i in range(n)]


def _render_chart_page(slide, s: Slide, tokens: dict) -> None:
    """图表页整幅：有结论就走「图表主区 + 结论侧栏」双栏，没有就让图表吃满整条正文带。
    此前图表卡在正文带里、结论压成一行小字贴在图下沿，左右各空出一大片——饼图尤其明显，
    它的直径被页高卡住，横向根本撑不到 12in（用户实评「图表页左右偏空」）。"""
    left, content_w, _ = content_box(slide)
    top, bottom = body_band(slide)
    notes = [b for b in s.bullets if b.strip()]
    chart_w = content_w
    if notes:
        rail_w = int(content_w * _RAIL_RATIO)
        chart_w = content_w - rail_w - CARD["col_gap"]
        _chart_rail(slide, notes, tokens, (left + chart_w + CARD["col_gap"], top, rail_w, bottom))
    _render_chart_body(slide, s.chart, tokens, (left, top, chart_w, bottom - top))


def _stat_card(slide, left, top, width, height, item: StatItem, tokens: dict) -> None:
    """关键数字大卡片：浅底色圆角矩形，大字号数字（主色）+ 一行说明（弱化灰）。
    comparison 版式的右栏专用——把最有冲击力的对比数字从要点文字里摘出来放大，
    比埋在项目符号列表里更有说服力。深色模板的数字取强调色：主色本身就是亮青绿，
    但在深底卡片上仍是最跳的一档。"""
    shape = rounded_rect(slide, left, top, width, height, tokens["tint"],
                         line_rgb=tokens["accent"])
    tf = shape.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = tf.margin_right = CARD["pad_x"]
    # 空串赋给 paragraph.text 不会产生 run，直接取 runs[0] 会 IndexError——编辑器「添加卡片」
    # 的初始值就是空串，用户没填就保存，之后每次导出都确定性崩（评审实测复现）。
    # 与 textbox 同一防御写法：没有 run 就补一个。渲染层只保证不崩，「不许留空」由入口校验负责。
    def _line(text: str, size: int, bold: bool, color: RGBColor, first: bool) -> None:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        p.text = text
        p.alignment = PP_ALIGN.CENTER
        p.line_spacing = LEAD["title"] if bold else LEAD["note"]
        if bold:
            p.space_after = Pt(8)      # 数字和说明之间要有一口气，挨着就分不出主次
        run = p.runs[0] if p.runs else p.add_run()
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color

    accent = tokens["accent"] if is_dark(tokens) else tokens["primary"]
    _line(item.value, TYPE["stat_value"], True, accent, first=True)
    _line(item.label, TYPE["note"], False, tokens["muted"], first=False)


def _render_comparison_body(slide, s: Slide, tokens: dict) -> None:
    """对比页主体：左栏要点（招标要求/传统方案的说明）+ 右栏 1-2 张数字大卡片（我方承诺/本方案的
    冲击力数字）。招标要求 vs 我方承诺、传统方案 vs 本方案这类内容用它，比堆一排项目符号更有说服力。"""
    left, content_w, _ = content_box(slide)
    gap = CARD["col_gap"]
    left_w = int(content_w * 0.56)
    right_left = left + left_w + gap
    right_w = content_w - left_w - gap
    bullets_box(slide, s.bullets, tokens, left=left, width=left_w)   # 左栏与要点页同款卡片
    top, bottom = body_band(slide)
    n = max(1, len(s.stats))
    card_h = int((bottom - top - (n - 1) * gap) / n)
    for i, item in enumerate(s.stats):
        _stat_card(slide, right_left, top + i * (card_h + gap), right_w, card_h, item, tokens)


def _render_content(slide, s: Slide, tokens: dict, n: int, total: int, kicker: str = "") -> None:
    """正文页：标题行 + 版式化主体（bullets/chart/comparison）+ 评分点角标（可空）+ 页码。"""
    title_row(slide, s.title, tokens, index=n, kicker=kicker)
    if s.layout == "chart" and s.chart:
        _render_chart_page(slide, s, tokens)
    elif s.layout == "comparison" and s.stats:
        _render_comparison_body(slide, s, tokens)
    elif s.bullets:
        bullets_box(slide, s.bullets, tokens)
    if s.scoring:
        _scoring_chip(slide, s.scoring, tokens)
    _page_number(slide, n, total, tokens)


_AI_NOTICE = "本内容由 AI 辅助生成，仅供参考，请人工复核后使用"


def _ai_notice(slide, tokens: dict) -> None:
    """结束页底部小字（spec326 算法备案）：强调条上方一行，10pt 弱化灰、居中，两条渲染路径共用。"""
    sw, sh = slide_size(slide)
    textbox(slide, sw * 0.1, sh - Inches(0.4), sw * 0.8, Inches(0.3),
            [_AI_NOTICE], size=TYPE["caption"], color=tokens["muted"], align=PP_ALIGN.CENTER)


def _render_end(slide, s: Slide, deck: DeckSpec, tokens: dict, n: int, total: int) -> None:
    """结束页：居中致谢标题（34pt 加粗主色）+ 项目名副标题（弱化灰）+ 底部强调条 + 页码 + AI 生成提示。"""
    sw, sh = slide_size(slide)
    title = s.title or "感谢聆听"
    rect(slide, (sw - Inches(0.9)) / 2, sh * 0.36, Inches(0.9), CARD["rule"], tokens["accent"])
    textbox(slide, sw * 0.12, sh * 0.40, sw * 0.76, Inches(1.1),
            [title], size=TYPE["section_title"], color=tokens["primary"], bold=True,
            align=PP_ALIGN.CENTER, line_spacing=LEAD["title"])
    if deck.title:
        textbox(slide, sw * 0.12, sh * 0.56, sw * 0.76, Inches(0.5),
                [deck.title], size=TYPE["section_lead"], color=tokens["muted"],
                align=PP_ALIGN.CENTER)
    accent_bar(slide, tokens["accent"])
    _page_number(slide, n, total, tokens)
    _ai_notice(slide, tokens)


def render_pptx(deck: DeckSpec, *, template: str | None = None,
                 master_bytes: bytes | None = None) -> bytes:
    """DeckSpec → .pptx 字节（确定性，无 LLM，§4.2.1 两段式的渲染段）。
    master_bytes=None（默认）→ 走 _render_blank。
    master_bytes 给定（企业自有 .pptx/.potx 母版）→ 尝试 _render_on_master，套用母版自身的
    主题/母版/版式/logo；母版加载或渲染过程任何异常（损坏文件、版式异常等）一律吞掉只记警告，
    回退 _render_blank——保证流水线里述标产物总能生成，不因客户母版问题整体失败。"""
    if master_bytes is not None:
        try:
            return _render_on_master(deck, template, master_bytes)
        except Exception:
            logger.warning("企业母版渲染失败，回退空白设计", exc_info=True)
    return _render_blank(deck, template)


def _render_blank(deck: DeckSpec, template: str | None) -> bytes:
    """空白设计路径：16:9，模板（blue/tech/gov）决定封面骨架/标题行/要点卡/配色；
    页码统计 content+end 页（封面不计分母/不显示页码）；口播稿写入备注页。"""
    tokens = tokens_for(template, deck.template)
    prs = Presentation()
    prs.slide_width, prs.slide_height = SLIDE_W, SLIDE_H
    blank = prs.slide_layouts[6]
    total = sum(1 for s in deck.slides if s.kind != "cover")
    n = sec = 0
    for s in deck.slides:
        slide = prs.slides.add_slide(blank)
        paint_surface(slide, tokens)      # 底色必须最先画，后续形状才叠在它上面
        if s.kind == "cover":
            render_cover(slide, s, tokens)
        elif s.kind == "section":
            n += 1
            sec += 1
            _render_section(slide, s, tokens, n, total, sec)
        elif s.kind == "end":
            n += 1
            _render_end(slide, s, deck, tokens, n, total)
        else:
            n += 1
            _render_content(slide, s, tokens, n, total, deck.title)
        if s.notes:
            slide.notes_slide.notes_text_frame.text = s.notes
    out = io.BytesIO()
    prs.save(out)
    return out.getvalue()


def _clear_slides(prs: Presentation) -> None:
    """删掉母版自带的示例页：既摘除 sldIdLst 引用也 drop 对应关系，让 slide part 在包里彻底
    不可达（只摘 sldIdLst 会留下孤儿 part，新增页可能复用同一 partname 导致 zip 内重名）。
    masters/layouts/theme 不在这条关系链上，不受影响。"""
    sld_id_lst = prs.slides._sldIdLst
    for sld_id in list(sld_id_lst):
        prs.part.drop_rel(sld_id.rId)
        sld_id_lst.remove(sld_id)


def _pick_content_layout(layouts: list) -> object:
    """内容版式启发式：优先名字含“Title and Content”/“content”；否则 index 1；
    否则名字含“blank”；否则 index 5/6；否则退回 index 0（layouts 非空由调用方保证）。"""
    for layout in layouts:
        if "content" in (layout.name or "").lower():
            return layout
    if len(layouts) > 1:
        return layouts[1]
    for layout in layouts:
        if "blank" in (layout.name or "").lower():
            return layout
    for idx in (5, 6):
        if len(layouts) > idx:
            return layouts[idx]
    return layouts[0]


def _pick_layouts(prs: Presentation) -> tuple:
    """从母版版式里选（标题版式, 内容版式）：标题版式优先 index 0（封面/结束页共用）。"""
    layouts = list(prs.slide_layouts)
    if not layouts:
        raise ValueError("母版没有可用版式")
    return layouts[0], _pick_content_layout(layouts)


_BODY_PLACEHOLDER_TYPES = (PP_PLACEHOLDER.BODY, PP_PLACEHOLDER.OBJECT, PP_PLACEHOLDER.SUBTITLE)


def _title_placeholder(slide):
    """母版标题占位符（idx=0/TITLE 类型），没有则 None。"""
    return slide.shapes.title


def _body_placeholder(slide):
    """母版正文/副标题占位符（非标题的 BODY/OBJECT/SUBTITLE 类型），没有则 None。"""
    for ph in slide.placeholders:
        if ph.placeholder_format.idx != 0 and ph.placeholder_format.type in _BODY_PLACEHOLDER_TYPES:
            return ph
    return None


def _fill_body_bullets(ph, bullets: list[str]) -> None:
    """把要点逐条写进母版正文占位符，一条一段（不加“• ”前缀，列表符号交给母版自身样式）。"""
    tf = ph.text_frame
    tf.clear()
    for i, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = bullet


def _render_cover_on_master(slide, s: Slide, tokens: dict) -> None:
    """封面页（母版路径）：有标题占位符就写标题（+ 正文/副标题占位符写 bullets）；
    没有占位符则整页退回空白设计的封面绘制（模板自己的封面骨架），版式换汤不换药。"""
    title_ph = _title_placeholder(slide)
    if title_ph is None:
        render_cover(slide, s, tokens)
        return
    title_ph.text_frame.text = s.title
    body_ph = _body_placeholder(slide)
    if body_ph is not None and s.bullets:
        _fill_body_bullets(body_ph, s.bullets)


def _render_content_on_master(slide, s: Slide, tokens: dict, n: int, total: int,
                              kicker: str = "") -> None:
    """正文页（母版路径）：标题优先落母版占位符，缺失则退回空白设计同款绘制。
    图表/对比版式的主体客户母版不会自带对应占位符，恒定自绘（同评分点角标/页码的既有做法）；
    只有 bullets 版式才尝试母版正文占位符，缺失同样退回空白设计同款绘制。"""
    title_ph = _title_placeholder(slide)
    if title_ph is not None:
        title_ph.text_frame.text = s.title
    else:
        title_row(slide, s.title, tokens, index=n, kicker=kicker)
    if s.layout == "chart" and s.chart:
        _render_chart_page(slide, s, tokens)
    elif s.layout == "comparison" and s.stats:
        _render_comparison_body(slide, s, tokens)
    elif s.bullets:
        body_ph = _body_placeholder(slide)
        if body_ph is not None:
            _fill_body_bullets(body_ph, s.bullets)
        else:
            bullets_box(slide, s.bullets, tokens)
    if s.scoring:
        _scoring_chip(slide, s.scoring, tokens)
    _page_number(slide, n, total, tokens)


def _render_end_on_master(slide, s: Slide, tokens: dict, n: int, total: int) -> None:
    """结束页（母版路径）：标题占位符写致谢语；缺失则退回空白设计的居中致谢绘制。
    页码恒定自绘；AI 生成提示恒定自绘（母版版式不会自带，两路径视觉一致）。"""
    title = s.title or "感谢聆听"
    title_ph = _title_placeholder(slide)
    if title_ph is not None:
        title_ph.text_frame.text = title
    else:
        sw, sh = slide_size(slide)
        textbox(slide, sw * 0.12, sh * 0.40, sw * 0.76, Inches(1.1),
                [title], size=TYPE["section_title"], color=tokens["primary"], bold=True,
                align=PP_ALIGN.CENTER, line_spacing=LEAD["title"])
    _page_number(slide, n, total, tokens)
    _ai_notice(slide, tokens)


def _render_on_master(deck: DeckSpec, template: str | None, master_bytes: bytes) -> bytes:
    """企业母版路径：加载客户 .pptx/.potx，清空母版自带示例页只留 masters/layouts/theme，
    再用母版自身版式承载我们的封面/正文/结束页（标题/正文占位符优先，缺失退回空白设计同款绘制）。
    不强制 16:9——沿用母版自身的页面尺寸（prs.slide_width/height 不改）。"""
    tokens = tokens_for(template, deck.template)
    prs = Presentation(io.BytesIO(master_bytes))
    _clear_slides(prs)
    title_layout, content_layout = _pick_layouts(prs)
    total = sum(1 for s in deck.slides if s.kind != "cover")
    n = sec = 0
    # 母版路径不铺底色：客户母版自带背景/底纹/logo，盖上去等于把人家的设计糊掉。
    for s in deck.slides:
        if s.kind == "cover":
            slide = prs.slides.add_slide(title_layout)
            _render_cover_on_master(slide, s, tokens)
        elif s.kind == "section":
            # 章节分隔页需要整页满色块自绘，客户母版的占位符不适合这种版式——沿用 title_layout
            # 只借它的页面尺寸/主题环境，视觉内容完全自绘（与空白设计路径一致）。
            n += 1
            sec += 1
            slide = prs.slides.add_slide(title_layout)
            _render_section(slide, s, tokens, n, total, sec)
        elif s.kind == "end":
            n += 1
            slide = prs.slides.add_slide(title_layout)
            _render_end_on_master(slide, s, tokens, n, total)
        else:
            n += 1
            slide = prs.slides.add_slide(content_layout)
            _render_content_on_master(slide, s, tokens, n, total, deck.title)
        if s.notes:
            slide.notes_slide.notes_text_frame.text = s.notes
    out = io.BytesIO()
    prs.save(out)
    return out.getvalue()
