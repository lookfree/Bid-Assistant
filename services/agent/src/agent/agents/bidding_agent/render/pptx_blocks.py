"""述标 PPT 的绘制积木：几何/图形原语 + 封面、标题行、要点卡三族版式画法 + 要点区密度版式。

模板之间的差异全部落在前三族上（配色、字号层级与卡片规格见 render/styles.py）。每族一个派发表，
渲染层只按 token 里的枚举取画法，**未知取值回退各自默认值**——加新模板只需在 styles.py
登记 token 并在这里补一个画法函数，deck 组装逻辑（render/pptx.py）一行都不用改。

密度族（_DENSITIES）是第四张派发表，但它**不由模板决定而由内容条数决定**：2 条和 6 条要点
用同一套竖排卡片，前者空得慌、后者挤成一团，而且每页长一个样。按条数换版式后，
同一套模板里页与页之间也有了节奏。
"""
from __future__ import annotations

import logging

from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

from agent.agents.bidding_agent.render.styles import (
    CARD, LEAD, TYPE, blend_toward, is_dark, on_primary,
)
from agent.agents.bidding_agent.schemas import Slide

logger = logging.getLogger(__name__)

SLIDE_W, SLIDE_H = Inches(13 + 1 / 3), Inches(7.5)  # 12192000 / 6858000 EMU：标准 16:9
MARGIN = Inches(0.7)


# ---- 几何：一律按幻灯片真实尺寸算，别取全局常量 ----

def slide_size(slide) -> tuple[int, int]:
    """该幻灯片所属演示文稿的真实页面尺寸（EMU）。
    企业母版路径**沿用客户自己的页面尺寸**（_render_on_master 不改 slide_width/height），
    死写 16:9 常量在 4:3 母版（10in 宽）上会把图表推出页面 2.63in——约四分之一被裁掉
    （评审实测复现）。取不到就退回 16:9 常量，与空白设计路径一致。"""
    try:
        pres = slide.part.package.presentation_part.presentation
        return pres.slide_width, pres.slide_height
    except Exception:      # noqa: BLE001 拿不到尺寸绝不能让整份述标渲染失败
        logger.warning("读取母版页面尺寸失败，按 16:9 常量绘制", exc_info=True)
        return SLIDE_W, SLIDE_H


def content_box(slide) -> tuple[int, int, int]:
    """(左边距, 正文可用宽, 页高)：所有版式统一用它算几何。"""
    w, h = slide_size(slide)
    return MARGIN, w - 2 * MARGIN, h


def body_band(slide) -> tuple[int, int]:
    """正文区上下沿：下沿按页高算，给底部的评分点角标/页码留出 1.2in（母版页高不同也不会压住）。"""
    _, _, sh = content_box(slide)
    return Inches(1.55), sh - Inches(1.2)


# ---- 图形/文字原语 ----

def rect(slide, left, top, width, height, fill_rgb, *, line_rgb=None, line_pt=None):
    """无框（或指定描边）的纯色矩形，封面色带/分隔线/强调条/标题小方块共用。"""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, int(left), int(top), int(width), int(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_rgb
    if line_rgb is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line_rgb
        shape.line.width = line_pt or CARD["hairline"]
    return shape


def rounded_rect(slide, left, top, width, height, fill_rgb, *, line_rgb=None, line_pt=None):
    """圆角矩形（要点卡片/数字卡片/角标共用），描边规则同 rect。
    圆角半径统一按 CARD["radius"] 的**绝对值**折算：圆角矩形的默认圆角是短边的 16.7%，
    于是 0.5in 高的角标是 0.08in 圆角、3in 高的数字卡是 0.5in 圆角——同一页上好几种圆角，
    一眼就看得出是"默认值堆出来的"。"""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                   int(left), int(top), int(width), int(height))
    short = max(1, min(int(width), int(height)))
    shape.adjustments[0] = min(0.5, CARD["radius"] / short)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_rgb
    if line_rgb is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line_rgb
        shape.line.width = line_pt or CARD["hairline"]
    return shape


def gradient_rect(slide, left, top, width, height, c1: RGBColor, c2: RGBColor, angle: float = 45.0):
    """渐变矩形。纯色大色块是"代码画的"最明显的特征——同一块面积换成渐变，观感立刻接近
    设计稿。python-pptx 原生支持渐变填充，产物仍是可编辑形状，不是图片。"""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                   int(left), int(top), int(width), int(height))
    shape.fill.gradient()
    stops = shape.fill.gradient_stops
    stops[0].color.rgb = c1
    stops[0].position = 0.0
    stops[1].color.rgb = c2
    stops[1].position = 1.0
    # 多于两个停靠点时把其余的并到末端，避免默认主题色混进来
    for extra in list(stops)[2:]:
        extra.color.rgb = c2
        extra.position = 1.0
    shape.fill.gradient_angle = angle
    shape.line.fill.background()
    return shape


def textbox(slide, left, top, width, height, lines, *, size, color, bold=False, align=None,
            line_spacing=None, space_after=None, tracking=None):
    """单个文本框，lines 逐行成段；每段单 run（够用且最稳）。

    **四边内边距一律清零**：文本框默认左右各留 0.1in 内边距，于是"放在 x 处的文字"实际画在
    x+0.1in——同一条左边线上的色块和标题就差了 0.1in，整页看着像没对齐（光学边距的第一现场）。
    位置由调用方给几何，需要留白就把留白算进几何里。
    tracking 是字距（pt）：小号标签字距放开一点才像标签，不放开就是一行小字。"""
    tf = slide.shapes.add_textbox(int(left), int(top), int(width), int(height)).text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        if align is not None:
            p.alignment = align
        if line_spacing is not None:
            p.line_spacing = line_spacing
        if space_after is not None:
            p.space_after = space_after
        run = p.runs[0] if p.runs else p.add_run()
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
        if tracking is not None:
            run.font._rPr.set("spc", str(int(tracking * 100)))   # 字距单位 = 1/100 pt
    return tf


def paint_surface(slide, tokens: dict) -> None:
    """深色模板铺满页底色（浅色模板不铺：空白版式本来就是白底，多画一层只会盖住母版元素）。
    必须最先画——后画的形状才会叠在它上面。"""
    if not is_dark(tokens):
        return
    sw, sh = slide_size(slide)
    rect(slide, 0, 0, sw, sh, tokens["bg"])


def accent_bar(slide, accent_rgb) -> None:
    """封面/结束页底部的通栏细强调条（0.06in），按真实页宽画——死写 16:9 常量会伸出窄母版之外。"""
    sw, sh = slide_size(slide)
    rect(slide, 0, sh - Inches(0.06), sw, Inches(0.06), accent_rgb)


def inset_frame(slide, inset: int, line_rgb: RGBColor, weight) -> None:
    """整页内缩一圈的细边框（四条线，不用描边矩形——描边矩形会连带一块填充要处理）。"""
    sw, sh = slide_size(slide)
    rect(slide, inset, inset, sw - 2 * inset, weight, line_rgb)
    rect(slide, inset, sh - inset - weight, sw - 2 * inset, weight, line_rgb)
    rect(slide, inset, inset, weight, sh - 2 * inset, line_rgb)
    rect(slide, sw - inset - weight, inset, weight, sh - 2 * inset, line_rgb)


# ---- 封面族：cover 开关 ----
# 三套封面各有一处**克制的视觉锚点**（几何构成，不用图片素材）：光有大色块和一行标题，
# 第一页就是"底色 + 字"，立不住。锚点只用线和块，投影/打印/黑白复印都不会糊掉。

def _cover_fullbleed(slide, s: Slide, tokens: dict) -> None:
    """满幅封面（商务提案）：整页主色渐变 + 顶部细线小标签，标题压在下半页、上方一道贯通细线分栏。
    视觉锚点＝**栏位节奏**：右半页四道等距竖线从页顶垂下、正好落在那条贯通线上，
    再加一道更粗的强调色竖线收尾——像一张排过版的封面，而不是一块刷满色的板子。"""
    sw, sh = slide_size(slide)
    fg = on_primary(tokens)
    deep = blend_toward(tokens["primary"], RGBColor(0, 0, 0), 0.42)
    gradient_rect(slide, 0, 0, sw, sh, deep, tokens["primary"], angle=90.0)
    rule_y = int(sh * 0.50)
    faint = blend_toward(fg, tokens["primary"], 0.74)
    for i in range(4):
        rect(slide, sw * (0.60 + 0.08 * i), 0, Pt(1), rule_y, faint)
    rect(slide, sw * 0.92, 0, Pt(4), int(sh * 0.34), tokens["accent"])
    left = sw * 0.07
    rect(slide, left, sh * 0.15, sw * 0.05, CARD["rule"], tokens["accent"])
    textbox(slide, left, sh * 0.185, sw * 0.6, Inches(0.4), ["述标演示"],
            size=TYPE["eyebrow"], color=blend_toward(fg, tokens["primary"], 0.28), tracking=1.8)
    rect(slide, left, rule_y, sw * 0.86, Pt(1), blend_toward(fg, tokens["primary"], 0.55))
    rect(slide, left, rule_y - Inches(0.06), Inches(0.16), Inches(0.16), tokens["accent"])
    textbox(slide, left, sh * 0.55, sw * 0.80, Inches(2.0), [s.title],
            size=TYPE["cover_title"], color=fg, bold=True, line_spacing=LEAD["title"])
    if s.bullets:
        textbox(slide, left, sh * 0.80, sw * 0.80, Inches(1.0), s.bullets,
                size=TYPE["cover_meta"], color=blend_toward(fg, tokens["primary"], 0.32),
                line_spacing=LEAD["note"], space_after=Pt(4))
    accent_bar(slide, tokens["accent"])


def _cover_split(slide, s: Slide, tokens: dict) -> None:
    """分栏封面（技术方案）：左栏主色实底压标题，右栏整片留白放副标题与投标人信息。
    视觉锚点＝右下角一片**细网格 + 一个裁切角标**：右栏原本是一大块空底色，
    网格给它一个"图纸"的地，裁切角标把版心的右下角钉住——留白仍是留白，但不再是空的。"""
    sw, sh = slide_size(slide)
    col = int(sw * 0.46)
    pad = int(sw * 0.055)
    fg = on_primary(tokens)
    rect(slide, 0, 0, col, sh, tokens["primary"])
    rect(slide, col, 0, Pt(3), sh, tokens["accent"])
    rect(slide, pad, sh * 0.30, Inches(0.7), CARD["rule"], fg)
    textbox(slide, pad, sh * 0.355, col - 2 * pad, Inches(2.4), [s.title],
            size=TYPE["cover_title"] - 6, color=fg, bold=True, line_spacing=LEAD["title"])
    textbox(slide, col + pad, sh * 0.30, sw - col - 2 * pad, Inches(0.5), ["述标演示"],
            size=TYPE["eyebrow"], color=tokens["accent"], tracking=1.8)
    if s.bullets:
        textbox(slide, col + pad, sh * 0.375, sw - col - 2 * pad, Inches(1.4), s.bullets,
                size=TYPE["cover_meta"], color=tokens["text"],
                line_spacing=LEAD["note"], space_after=Pt(4))
    _blueprint_patch(slide, col + pad, int(sh * 0.60), sw - col - 2 * pad, int(sh * 0.28), tokens)


_DOT = Inches(0.045)
_DOT_COLS, _DOT_ROWS = 7, 4


def _blueprint_patch(slide, left, top, width, height, tokens: dict) -> None:
    """点阵 + 裁切角标。用点不用整根线：整根线画出来就是一张**空表格**（实测第一版如此，
    像忘了填内容），点阵只是一层纹理，给空面板一个"排过版"的地。"""
    dot = blend_toward(tokens["bg"], tokens["text"], 0.34)
    for i in range(_DOT_COLS):
        for j in range(_DOT_ROWS):
            rect(slide, left + (width - _DOT) * i / (_DOT_COLS - 1),
                 top + (height - _DOT) * j / (_DOT_ROWS - 1), _DOT, _DOT, dot)
    mark = int(min(width, height) * 0.3)
    rect(slide, left + width - mark, top + height + Inches(0.22), mark, Pt(2.5), tokens["accent"])
    rect(slide, left + width - Pt(2.5), top + height - mark + Inches(0.22), Pt(2.5), mark,
         tokens["accent"])


def _cover_banner(slide, s: Slide, tokens: dict) -> None:
    """通栏横幅封面（党政庄重）：一条横贯页面的主色带压住标题，上沿一道烫金细线。
    视觉锚点＝**内缩一圈的烫金细边框 + 横幅左端的竖向绶带块 + 横幅下沿的第二道细线**：
    边框与双线是正式文书的处理，绶带块给横幅一个起点，整页从"白纸一条红带"变成有仪式感的封面。"""
    sw, sh = slide_size(slide)
    top, band_h = int(sh * 0.30), int(sh * 0.30)
    left = int(sw * 0.09)
    inset_frame(slide, int(min(sw, sh) * 0.045), blend_toward(tokens["accent"], tokens["white"], 0.4),
                Pt(0.75))
    rect(slide, 0, top - Pt(5), sw, Pt(5), tokens["accent"])
    rect(slide, 0, top, sw, band_h, tokens["primary"])
    rect(slide, 0, top + band_h, sw, Pt(2), blend_toward(tokens["accent"], tokens["white"], 0.25))
    rect(slide, int(sw * 0.055), top - Inches(0.22), Inches(0.14), band_h + Inches(0.44),
         tokens["accent"])
    textbox(slide, left, top - Inches(0.7), sw - 2 * left, Inches(0.4), ["述标演示"],
            size=TYPE["eyebrow"], color=tokens["muted"], tracking=1.8)
    textbox(slide, left, top + band_h * 0.22, sw - 2 * left, band_h * 0.66, [s.title],
            size=TYPE["cover_title"] - 6, color=on_primary(tokens), bold=True,
            line_spacing=LEAD["title"])
    if s.bullets:
        textbox(slide, left, top + band_h + Inches(0.45), sw - 2 * left, Inches(1.0),
                s.bullets, size=TYPE["cover_meta"], color=tokens["text"],
                line_spacing=LEAD["note"], space_after=Pt(4))
    accent_bar(slide, tokens["accent"])


_COVERS = {"fullbleed": _cover_fullbleed, "split": _cover_split, "banner": _cover_banner}


def render_cover(slide, s: Slide, tokens: dict) -> None:
    """封面页：按 cover 开关取画法，未知取值回退满幅封面（最通用的一款）。"""
    _COVERS.get(tokens.get("cover"), _cover_fullbleed)(slide, s, tokens)


# ---- 标题行族：header 开关 ----

def _title_overline(slide, title: str, tokens: dict, index: int | None, kicker: str) -> None:
    """标题行（商务提案）：细线 + 小标签（overline）压在标题上方，标题下不再画通栏线——
    小标签走项目名（running head），评委翻到任意一页都知道这是哪个项目的述标。"""
    left, content_w, _ = content_box(slide)
    rect(slide, left, Inches(0.44), Inches(0.5), CARD["rule"], tokens["accent"])
    textbox(slide, left + Inches(0.66), Inches(0.36), content_w - Inches(0.66), Inches(0.32),
            [kicker or "述标演示"], size=TYPE["eyebrow"], color=tokens["accent"], tracking=1.2)
    textbox(slide, left, Inches(0.68), content_w, Inches(0.7), [title], size=TYPE["page_title"],
            color=tokens["text"], bold=True, line_spacing=LEAD["title"])


def _title_numeral(slide, title: str, tokens: dict, index: int | None, kicker: str) -> None:
    """标题行（技术方案）：大号章节数字起标题，除此之外什么都不画——正文首行的发丝线本身
    就是分隔（再画一道标题线会和它挨成两条平行线，实测像画歪了而不像设计）。"""
    left, content_w, _ = content_box(slide)
    num_w = Inches(1.1)
    textbox(slide, left, Inches(0.42), num_w, Inches(0.9),
            [f"{index:02d}" if index else "—"], size=TYPE["page_title"] + 6,
            color=tokens["accent"], bold=True, line_spacing=LEAD["title"])
    textbox(slide, left + num_w, Inches(0.52), content_w - num_w, Inches(0.7), [title],
            size=TYPE["page_title"], color=tokens["text"], bold=True, line_spacing=LEAD["title"])


def _title_corner(slide, title: str, tokens: dict, index: int | None, kicker: str) -> None:
    """标题行（党政庄重）：贴左边缘的主色角标（内嵌页序）+ 标题 + 标题下半幅烫金短线。"""
    left, content_w, _ = content_box(slide)
    tab_w = left + Inches(0.36)
    rect(slide, 0, Inches(0.42), tab_w, Inches(0.78), tokens["primary"])
    textbox(slide, 0, Inches(0.6), tab_w - Inches(0.14), Inches(0.4),
            [f"{index:02d}" if index else "标"], size=TYPE["eyebrow"] + 2,
            color=on_primary(tokens), bold=True, align=PP_ALIGN.RIGHT)
    textbox(slide, tab_w + Inches(0.26), Inches(0.44), content_w - Inches(0.26), Inches(0.66),
            [title], size=TYPE["page_title"], color=tokens["text"], bold=True,
            line_spacing=LEAD["title"])
    rect(slide, tab_w + Inches(0.26), Inches(1.2), content_w * 0.42, CARD["rule"], tokens["accent"])


_TITLE_ROWS = {"overline": _title_overline, "numeral": _title_numeral, "corner": _title_corner}


def title_row(slide, title: str, tokens: dict, *, index: int | None = None, kicker: str = "") -> None:
    """正文页标题行：按 header 开关取画法，未知取值回退 overline。
    模板之间的区分主要就体现在这里，别在调用方按模板名分支。"""
    _TITLE_ROWS.get(tokens.get("header"), _title_overline)(slide, title, tokens, index, kicker)


# ---- 要点卡族：card 开关 ----

_BADGE = Inches(0.34)


def _card_text(shape, text: str, tokens: dict, *, left_margin: int, size: int) -> None:
    """卡片内的要点文字：垂直居中、按 left_margin 给序号让位。"""
    tf = shape.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left, tf.margin_right = left_margin, CARD["pad_x"]
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = PP_ALIGN.LEFT
    p.line_spacing = LEAD["body"]
    run = p.runs[0] if p.runs else p.add_run()
    run.font.size = Pt(size)
    run.font.color.rgb = tokens["text"]


def _badge(slide, left, top, idx: int, tokens: dict) -> None:
    """强调色圆角序号徽章（白字加粗）。"""
    badge = rounded_rect(slide, left, top, _BADGE, _BADGE, tokens["accent"])
    tf = badge.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.text = str(idx)
    p.alignment = PP_ALIGN.CENTER
    run = p.runs[0] if p.runs else p.add_run()
    run.font.size = Pt(TYPE["eyebrow"] + 1)
    run.font.bold = True
    run.font.color.rgb = tokens["white"]


def _card_numbered(slide, left, top, width, height, idx: int, text: str, tokens: dict,
                   *, size: int) -> None:
    """要点卡（商务提案）：浅底圆角卡 + 左侧大号序号。序号是正文的两倍大，
    评委扫一眼就知道这页有几条、看到第几条；正文放大时序号跟着放大，层级不会被拉平。"""
    card = rounded_rect(slide, left, top, width, height, tokens["tint"])
    num_h = Inches(0.04) * size
    textbox(slide, left + CARD["pad_x"], top + (height - num_h) / 2, Inches(0.82), num_h,
            [str(idx)], size=size * 2, color=tokens["accent"], bold=True, align=PP_ALIGN.CENTER,
            line_spacing=LEAD["title"])
    _card_text(card, text, tokens, left_margin=CARD["pad_x"] + Inches(0.86), size=size)


def _card_hairline(slide, left, top, width, height, idx: int, text: str, tokens: dict,
                   *, size: int) -> None:
    """要点行（技术方案）：不填底色，只在行首压一道发丝线 + 小号序号——留白即版式，
    深底上大面积浅色块会盖过图表，这套模板宁可什么都不填。"""
    rect(slide, left, top, width, CARD["hairline"], tokens["muted"])
    # 序号与正文都按整行垂直居中：序号盒子若只给一行高，会浮在正文上方半行，看着像没对齐
    num = textbox(slide, left, top + Inches(0.12), Inches(0.6), height - Inches(0.24),
                  [f"{idx:02d}"], size=TYPE["eyebrow"], color=tokens["accent"], bold=True)
    num.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf = textbox(slide, left + Inches(0.68), top + Inches(0.12), width - Inches(0.78),
                 height - Inches(0.24), [text], size=size, color=tokens["text"],
                 line_spacing=LEAD["body"])
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE


def _card_elevated(slide, left, top, width, height, idx: int, text: str, tokens: dict,
                   *, size: int) -> None:
    """要点卡（党政庄重）：卡片下压一层偏移色块当轻阴影 + 烫金描边 + 序号徽章。
    python-pptx 没有阴影 API，用一层偏移块近似——产物仍是可编辑的原生形状，不是图片。"""
    off = Inches(0.05)
    rounded_rect(slide, left + off, top + off, width, height,
                 blend_toward(tokens["tint"], tokens["text"], 0.18))
    card = rounded_rect(slide, left, top, width, height, tokens["tint"],
                        line_rgb=tokens["accent"])
    _badge(slide, left + CARD["pad_x"] - Inches(0.06), top + (height - _BADGE) / 2, idx, tokens)
    _card_text(card, text, tokens, left_margin=CARD["pad_x"] + Inches(0.36), size=size)


_CARDS = {"numbered": _card_numbered, "hairline": _card_hairline, "elevated": _card_elevated}


# ---- 要点区密度族：按条目数换版式 ----
# (列数, 单卡最大高度, 正文字号角色)。用户实评：2 条和 6 条长同一个样，前者空、后者挤。
_DENSITIES = {
    "lead":    (1, Inches(1.95), "lead"),      # ≤3 条：大卡纵向，字也跟着放大
    "stack":   (1, Inches(1.30), "body"),      # 4–6 条但栏很窄（对比页左栏）：仍走单列
    "duo":     (2, Inches(1.85), "body"),      # 4–6 条：双栏，一行两张
    "compact": (1, Inches(0.88), "compact"),   # >6 条：紧凑行，别把一页撑爆
}


def _density_for(n: int, wide: bool) -> str:
    """条目数（和这一栏够不够宽）→ 密度名。窄栏一律不分双栏：对比页左栏只有 56% 宽，
    再切两半每张卡放不下一句话。"""
    if n <= 3:
        return "lead"
    if n > 6:
        return "compact"
    return "duo" if wide else "stack"


def _grid_cells(n: int, cols: int, max_h: int, box: tuple[int, int, int, int]) -> list[tuple]:
    """n 张卡在 (left, top, width, bottom) 里的逐张矩形：按列铺、行内等宽、整体垂直居中。
    垂直居中是为了消灭「内容挤在顶部、下面 60% 全白」——用户原话「太素」，一半原因就是这片空白。"""
    left, top, width, bottom = box
    rows = -(-n // cols)
    gap, col_gap = CARD["gap"], CARD["col_gap"]
    area = bottom - top
    h = min(max_h, int((area - (rows - 1) * gap) / rows))
    y0 = top + int((area - (rows * h + (rows - 1) * gap)) / 2)
    col_w = int((width - (cols - 1) * col_gap) / cols)
    cells = []
    for i in range(n):
        row, col = divmod(i, cols)
        cells.append((left + col * (col_w + col_gap), y0 + row * (h + gap), col_w, h))
    if cols > 1 and n % cols:
        # 双栏里落单的最后一张横跨整行：留一个半格空洞比排满更难看
        x, y, w, h = cells[-1]
        cells[-1] = (left, y, width, h)
    return cells


def bullets_box(slide, bullets: list[str], tokens: dict, *, left=None, width=None) -> None:
    """要点区：按 card 开关逐条渲卡片（未知取值回退 numbered），版式按条目数选密度，
    整体垂直居中（对比页左栏传 left/width 复用同款卡片，窄栏自动退回单列）。"""
    items = [b for b in bullets if b.strip()]
    if not items:
        return
    box_left, box_w, _ = content_box(slide)
    left = box_left if left is None else left
    width = box_w if width is None else width
    top, bottom = body_band(slide)
    cols, max_h, role = _DENSITIES[_density_for(len(items), width >= box_w * 0.8)]
    painter = _CARDS.get(tokens.get("card"), _card_numbered)
    for i, (x, y, w, h) in enumerate(_grid_cells(len(items), cols, max_h, (left, top, width, bottom))):
        painter(slide, x, y, w, h, i + 1, items[i], tokens, size=TYPE[role])
