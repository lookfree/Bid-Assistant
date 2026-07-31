"""POST /render/deck：从存库 deck 确定性重渲述标 .pptx。

存在的理由（生产缺陷，2026-07-30）：述标页的「导出」只是取预签名 URL 直下 MinIO 里的旧对象，
既不重渲也不作废——用户在编辑器里改完 deck 再导出，拿到的仍是编辑前那份 PPT，可能就这么带去投标
（标书正文侧 7-28 已修过同样的病）。而 export 步对「用户自己上传标书」的 review-kind 项目一律拒绝，
那条路连 pptx key 都没人存过（present 步的 artifacts 被 executor 丢掉，只留 result），
所以不能靠「导出前先跑 export 步」来补。

统一解法就是这个接口：到渲染这一层，标书来自流水线正文还是用户上传已经无所谓了——都只是一个 deck。
无 LLM、不进 thread、不涉计费，与 /render/checklist 同范式（agent 落 MinIO 返 key，App 负责预签名）。
"""
import io

import pytest
from pptx import Presentation

from agent.routes import deck as deck_mod
from agent.routes.deck import DeckRenderBody, render_deck

_DECK = {
    "title": "零信任平台述标",
    "duration": 15,
    "template": "blue",
    "slides": [
        {"id": "s0", "title": "封面", "kind": "cover", "bullets": []},
        {"id": "s1", "title": "质保承诺", "kind": "content", "bullets": ["7×24 响应"],
         "scoring": "售后服务 10 分"},
    ],
}


class _Storage:
    def __init__(self):
        self.calls: list[tuple[str, bytes, str]] = []

    async def put_bytes(self, key, data, content_type=None):
        self.calls.append((key, data, content_type))


async def test_render_deck_uploads_to_the_thread_key_and_returns_it(monkeypatch):
    """落 artifacts/<thread_id>/present.pptx——与 present 节点同一个确定性 key，重渲即覆盖，
    下载侧不必知道这次渲染是谁触发的。"""
    store = _Storage()
    monkeypatch.setattr(deck_mod, "storage", store)
    res = await render_deck(DeckRenderBody(thread_id="proj-abc", deck=_DECK))
    assert len(store.calls) == 1
    key, data, ct = store.calls[0]
    assert key == "artifacts/proj-abc/present.pptx"
    assert res == {"key": key}
    assert "presentationml" in ct
    prs = Presentation(io.BytesIO(data))
    assert len(prs.slides) == 2


async def test_render_deck_rejects_a_thread_id_that_escapes_the_prefix(monkeypatch):
    """thread_id 直接拼进对象 key，必须挡住路径穿越——否则调用方一个 '../' 就能覆盖别的对象。"""
    store = _Storage()
    monkeypatch.setattr(deck_mod, "storage", store)
    for bad in ("../evil", "a/b", "", "proj/../x"):
        with pytest.raises(ValueError):
            await render_deck(DeckRenderBody(thread_id=bad, deck=_DECK))
    assert store.calls == []


async def test_render_deck_reflects_edits_made_after_generation(monkeypatch):
    """本接口存在的全部意义：改过的 deck 必须体现在产物里。"""
    store = _Storage()
    monkeypatch.setattr(deck_mod, "storage", store)
    edited = {**_DECK, "slides": [
        {"id": "s0", "title": "封面", "kind": "cover", "bullets": []},
        {"id": "s1", "title": "质保承诺", "kind": "content",
         "bullets": ["改成 4 小时到场"], "scoring": "售后服务 10 分"},
    ]}
    await render_deck(DeckRenderBody(thread_id="proj-abc", deck=edited))
    prs = Presentation(io.BytesIO(store.calls[0][1]))
    texts = [sh.text_frame.text for sh in prs.slides[1].shapes if sh.has_text_frame]
    assert any("4 小时到场" in t for t in texts)


async def test_render_deck_survives_a_broken_enterprise_master(monkeypatch):
    """企业母版取不到/损坏时回退空白设计，绝不让导出失败——与 present/export 既有兜底一致。"""
    store = _Storage()
    monkeypatch.setattr(deck_mod, "storage", store)

    async def _boom(key):
        return None

    monkeypatch.setattr(deck_mod, "fetch_master_bytes", _boom)
    res = await render_deck(DeckRenderBody(thread_id="proj-abc", deck=_DECK,
                                           enterprise_template_key="missing/master.pptx"))
    assert res["key"] == "artifacts/proj-abc/present.pptx"
    assert len(store.calls) == 1
