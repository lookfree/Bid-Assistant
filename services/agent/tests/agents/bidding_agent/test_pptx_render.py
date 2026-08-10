import io
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_SHAPE_TYPE
from pptx.util import Emu, Inches
from agent.agents.bidding_agent.schemas import DeckSpec
from agent.agents.bidding_agent.render.pptx import render_pptx
from agent.agents.bidding_agent.render.styles import CARD, TYPE
from agent.agents.bidding_agent.render.styles import TEMPLATE_TOKENS as _TEMPLATE_TOKENS


def _deck():
    return DeckSpec(title="述标", slides=[
        {"id": "s0", "title": "封面", "bullets": ["客户：某局", "时长 15 分钟"], "kind": "cover"},
        {"id": "s1", "title": "运维体系", "bullets": ["7×24 值守", "分级 SLA", "故障 30 分钟响应"],
         "scoring": "技术方案 50 分", "notes": "讲稿…", "kind": "content"},
        {"id": "s2", "title": "感谢聆听", "kind": "end"},
    ])


def test_render_pptx_produces_valid_deck():
    data = render_pptx(_deck())
    assert data[:2] == b"PK"                       # .pptx 是 zip
    prs = Presentation(io.BytesIO(data))
    assert len(prs.slides) == 3
    assert prs.slides[1].notes_slide.notes_text_frame.text == "讲稿…"


def test_slide_is_16_by_9():
    data = render_pptx(_deck())
    prs = Presentation(io.BytesIO(data))
    assert prs.slide_width == Emu(12192000)   # Inches(13.333)
    assert prs.slide_height == Emu(6858000)   # Inches(7.5)


def _solid_rgb(shape):
    """形状的纯色填充色；渐变/无填充/非图形一律给 None（直接取 fore_color 会抛错）。"""
    try:
        return shape.fill.fore_color.rgb
    except (AttributeError, TypeError, ValueError):
        return None


def test_blue_cover_is_a_full_bleed_gradient_with_the_title_in_the_lower_half():
    """商务提案的封面是满幅色块（走渐变——纯色大块最像"代码画的"），标题压在下半页，
    投标人信息跟在标题下方：整页没有一块无意义的白。"""
    from pptx.enum.dml import MSO_FILL
    prs = Presentation(io.BytesIO(render_pptx(_deck(), template="blue")))
    cover = prs.slides[0]
    grads = [sh for sh in cover.shapes
             if sh.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE and sh.fill.type == MSO_FILL.GRADIENT]
    assert grads, "封面色块没有用渐变"
    full = grads[0]
    assert full.width == prs.slide_width and full.height == prs.slide_height, "封面色块没有铺满整页"
    assert full.fill.gradient_stops[1].color.rgb == _TEMPLATE_TOKENS["blue"]["primary"]

    title_box = next(sh for sh in cover.shapes if sh.has_text_frame and sh.text_frame.text == "封面")
    run = title_box.text_frame.paragraphs[0].runs[0]
    assert run.font.size.pt == TYPE["cover_title"]
    assert run.font.bold is True
    assert run.font.color.rgb == RGBColor(0xFF, 0xFF, 0xFF)
    assert title_box.top > prs.slide_height * 0.5, "满幅封面的标题要压在下半页"
    meta = next(sh for sh in cover.shapes if sh.has_text_frame and "客户：某局" in sh.text_frame.text)
    assert meta.top > title_box.top


def test_tech_cover_is_split_into_two_columns():
    """技术方案的封面是左右分栏：左栏一块主色实底压标题，右栏整片留白放投标人信息。"""
    prs = Presentation(io.BytesIO(render_pptx(_deck(), template="tech")))
    cover = prs.slides[0]
    col = next(sh for sh in cover.shapes
               if sh.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE and sh.height == prs.slide_height
               and _solid_rgb(sh) == _TEMPLATE_TOKENS["tech"]["primary"])
    assert prs.slide_width * 0.3 < col.width < prs.slide_width * 0.6, "左栏不是分栏宽度"
    title_box = next(sh for sh in cover.shapes if sh.has_text_frame and sh.text_frame.text == "封面")
    assert title_box.left + title_box.width <= col.width, "标题要落在左栏里"
    meta = next(sh for sh in cover.shapes if sh.has_text_frame and "客户：某局" in sh.text_frame.text)
    assert meta.left >= col.width, "投标人信息要落在右栏留白里"


def test_gov_cover_is_a_full_width_banner():
    """党政庄重的封面是通栏横幅：一条横贯页面的主色带压住标题，上下留白。"""
    prs = Presentation(io.BytesIO(render_pptx(_deck(), template="gov")))
    cover = prs.slides[0]
    banner = next(sh for sh in cover.shapes
                  if sh.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE and sh.width == prs.slide_width
                  and _solid_rgb(sh) == _TEMPLATE_TOKENS["gov"]["primary"])
    assert banner.top > 0, "横幅不该贴着页顶——上方要留白"
    assert banner.height < prs.slide_height * 0.5, "横幅是一条带子，不是满幅色块"
    title_box = next(sh for sh in cover.shapes if sh.has_text_frame and sh.text_frame.text == "封面")
    assert banner.top <= title_box.top <= banner.top + banner.height, "标题要压在横幅里"


def test_content_slide_bullets_and_scoring_chip():
    """要点从「单文本框 + • 前缀」改成逐条编号卡片（用户反馈整页太素、下半页全白）：
    一条要点一张卡，卡上有序号徽章。全部是原生形状，仍可在 PowerPoint 里逐个改。
    3 条要点走 lead 密度（大卡），正文字号随之放大到 TYPE["lead"]。"""
    data = render_pptx(_deck())
    prs = Presentation(io.BytesIO(data))
    content = prs.slides[1]
    cards = [sh for sh in content.shapes
             if sh.has_text_frame and sh.text_frame.text in
             ("7×24 值守", "分级 SLA", "故障 30 分钟响应")]
    assert len(cards) == 3
    for card in cards:
        assert card.text_frame.paragraphs[0].runs[0].font.size.pt == TYPE["lead"]
    badges = [sh for sh in content.shapes
              if sh.has_text_frame and sh.text_frame.text in ("1", "2", "3")]
    assert len(badges) == 3, "每张卡片要有序号徽章"
    chip = next(sh for sh in content.shapes
                if sh.has_text_frame and "评分点｜" in sh.text_frame.text)
    assert chip.text_frame.text == "评分点｜技术方案 50 分"
    page_no = next(sh for sh in content.shapes
                   if sh.has_text_frame and "/" in sh.text_frame.text and sh is not chip)
    assert page_no.text_frame.text == "1 / 2"


def test_end_slide_has_thank_you_and_page_number():
    data = render_pptx(_deck())
    prs = Presentation(io.BytesIO(data))
    end = prs.slides[2]
    texts = [sh.text_frame.text for sh in end.shapes if sh.has_text_frame]
    assert "感谢聆听" in texts
    assert "2 / 2" in texts


def _tiny_master(width_in: float = 10.0, height_in: float = 7.5) -> bytes:
    """构造一个自带 1 张示例页的迷你母版（复用 python-pptx 内置模板的 layouts/theme）。"""
    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(width_in), Inches(height_in)
    prs.slides.add_slide(prs.slide_layouts[0])   # 母版自带的示例页，渲染时应被清空
    out = io.BytesIO()
    prs.save(out)
    return out.getvalue()


def test_render_on_master_removes_example_slide_and_keeps_master_size():
    data = render_pptx(_deck(), master_bytes=_tiny_master())
    prs = Presentation(io.BytesIO(data))
    assert len(prs.slides) == 3                 # 示例页被删，只剩我们的 3 页
    assert prs.slide_width == Inches(10.0)       # 母版自身尺寸保留，不强制 16:9
    assert prs.slide_height == Inches(7.5)


def test_render_on_master_populates_titles_notes_and_chip():
    data = render_pptx(_deck(), master_bytes=_tiny_master())
    prs = Presentation(io.BytesIO(data))
    all_texts = [sh.text_frame.text for sl in prs.slides for sh in sl.shapes if sh.has_text_frame]
    assert {"封面", "运维体系", "感谢聆听"} <= set(all_texts)
    assert prs.slides[1].notes_slide.notes_text_frame.text == "讲稿…"
    chip = next(sh for sl in prs.slides for sh in sl.shapes
                if sh.has_text_frame and "评分点｜" in sh.text_frame.text)
    assert chip.text_frame.text == "评分点｜技术方案 50 分"


def _all_texts(prs: Presentation) -> list[str]:
    return [sh.text_frame.text for sl in prs.slides for sh in sl.shapes if sh.has_text_frame]


def test_end_slide_has_ai_notice_blank_path():
    """spec326 算法备案：结束页（空白设计路径）底部含 AI 生成提示短版文案（逐字不可改）。"""
    data = render_pptx(_deck())
    prs = Presentation(io.BytesIO(data))
    assert "本内容由 AI 辅助生成，仅供参考，请人工复核后使用" in _all_texts(prs)


def test_end_slide_has_ai_notice_master_path():
    """spec326：结束页（企业母版路径）同样含 AI 生成提示短版文案，两路径视觉一致。"""
    data = render_pptx(_deck(), master_bytes=_tiny_master())
    prs = Presentation(io.BytesIO(data))
    assert "本内容由 AI 辅助生成，仅供参考，请人工复核后使用" in _all_texts(prs)


def test_render_on_master_malformed_bytes_falls_back_to_blank():
    data = render_pptx(_deck(), master_bytes=b"not a pptx")
    assert data[:2] == b"PK"
    prs = Presentation(io.BytesIO(data))
    assert len(prs.slides) == 3
    assert prs.slide_width == Emu(12192000)      # 回退空白设计：强制 16:9


# ---- 述标结构性升级：章节分隔页 / 图表页 / 对比页（三种新版式） ----

def _rich_deck():
    return DeckSpec(title="述标", slides=[
        {"id": "s0", "title": "封面", "bullets": ["客户：某局"], "kind": "cover"},
        {"id": "sec", "title": "技术方案", "bullets": ["核心能力与差异化优势"], "kind": "section"},
        {"id": "s1", "title": "团队构成", "kind": "content", "layout": "chart", "scoring": "团队 20 分",
         "bullets": ["60% 为中级及以上职称"],
         "chart": {"type": "pie", "categories": ["高级", "中级", "初级"],
                   "series": [{"name": "人数", "values": [3, 6, 4]}]}},
        {"id": "s2", "title": "业绩对比", "kind": "content", "layout": "comparison", "scoring": "业绩 15 分",
         "bullets": ["近三年同类项目 5 个", "合同额年增长 30%"],
         "stats": [{"value": "72 小时", "label": "较招标要求提前完成"},
                   {"value": "0 起", "label": "质量投诉记录"}]},
        {"id": "s3", "title": "结语", "kind": "end"},
    ])


def test_section_slide_is_full_color_with_centered_title():
    """章节分隔页：满屏主色块 + 居中大标题，不挂评分点角标（它是过渡页，不对应具体得分点）。"""
    data = render_pptx(_rich_deck())
    prs = Presentation(io.BytesIO(data))
    section = prs.slides[1]
    rects = [sh for sh in section.shapes if sh.shape_type == MSO_SHAPE.RECTANGLE]
    assert any(r.width == prs.slide_width and r.height == prs.slide_height for r in rects)
    texts = [sh.text_frame.text for sh in section.shapes if sh.has_text_frame]
    assert "技术方案" in texts
    assert not any("评分点｜" in t for t in texts)


def test_chart_slide_renders_a_real_editable_chart_not_an_image():
    """图表页：真实 PowerPoint 图表对象（python-pptx add_chart），评委能在 PPT 里直接编辑数值——
    这正是相对"糊一张图片上去"的核心差异，也是本次结构升级要验证的主张。"""
    data = render_pptx(_rich_deck())
    prs = Presentation(io.BytesIO(data))
    chart_slide = prs.slides[2]
    chart_shapes = [sh for sh in chart_slide.shapes if sh.has_chart]
    assert len(chart_shapes) == 1
    chart = chart_shapes[0].chart
    assert list(chart.plots[0].categories) == ["高级", "中级", "初级"]
    assert [s.name for s in chart.series] == ["人数"]
    assert list(chart.series[0].values) == [3, 6, 4]
    # 评分点角标仍在——图表版式不能因为换了主体就丢了述标的核心标注
    texts = [sh.text_frame.text for sh in chart_slide.shapes if sh.has_text_frame]
    assert any("评分点｜团队 20 分" in t for t in texts)


def test_comparison_slide_has_left_bullets_and_right_stat_cards():
    """对比页：左栏要点 + 右栏 1-2 张数字大卡片，两栏都要有——这是招标要求 vs 我方承诺、
    传统方案 vs 本方案这类内容该用的版式，比堆一排项目符号更有说服力。"""
    data = render_pptx(_rich_deck())
    prs = Presentation(io.BytesIO(data))
    cmp_slide = prs.slides[3]
    texts = [sh.text_frame.text for sh in cmp_slide.shapes if sh.has_text_frame]
    assert any("近三年同类项目 5 个" in t for t in texts)
    # shape_type 对所有自选图形都返回 AUTO_SHAPE，具体形状要看 auto_shape_type
    # （auto_shape_type 对非自选图形直接抛 ValueError，不是返回 None，得先判 shape_type）
    cards = [sh for sh in cmp_slide.shapes
             if sh.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE and sh.auto_shape_type == MSO_SHAPE.ROUNDED_RECTANGLE]
    card_texts = {c.text_frame.text for c in cards}
    assert any("72 小时" in t and "较招标要求提前完成" in t for t in card_texts)
    assert any("0 起" in t and "质量投诉记录" in t for t in card_texts)


def test_new_layouts_render_on_enterprise_master_too():
    """企业母版路径：章节分隔页/图表页/对比页都是自绘主体（客户模板不会自带这些占位符），
    不因为换了母版就整段消失或报错——同评分点角标/页码「母版不自带、恒定自绘」的既有约定。"""
    data = render_pptx(_rich_deck(), master_bytes=_tiny_master())
    prs = Presentation(io.BytesIO(data))
    assert len(prs.slides) == 5
    assert any(sh.has_chart for sl in prs.slides for sh in sl.shapes)
    assert any(sh.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE and sh.auto_shape_type == MSO_SHAPE.ROUNDED_RECTANGLE
               for sl in prs.slides for sh in sl.shapes)


# ---- 评审回归：渲染层绝不能因为"内容形状特殊"整份崩掉 ----

def test_blank_stat_card_does_not_crash_the_render():
    """空串赋给 paragraph.text 不产生 run，取 runs[0] 会 IndexError → 付费导出确定性失败且重试
    无效（评审实测复现）。入口现在拒绝空卡片（StatItem 与 App PATCH 都要求非空），但**渲染层
    仍须自己防御**：库里可能已有旧数据，且导出每次都会重渲存量 deck——渲染层崩掉就再也导不出来。
    因此用 model_construct 绕过校验直接构造"历史脏数据"形状。"""
    from agent.agents.bidding_agent.schemas import StatItem
    blank = StatItem.model_construct(value="", label="")
    deck = DeckSpec(title="述标", slides=[
        {"id": "s1", "title": "对比", "kind": "content", "layout": "comparison", "bullets": ["要点"]},
    ])
    deck.slides[0].stats = [blank]
    prs = Presentation(io.BytesIO(render_pptx(deck)))
    assert len(prs.slides) == 1


def test_many_category_pie_does_not_crash_the_render():
    """报价构成/岗位分布这类饼图 7 个类别很常见。配色的混合比例不封顶会算出 >255 的通道值，
    RGBColor 直接抛 ValueError → 整个 present run 在模型已经花完钱之后失败（评审实测复现）。"""
    deck = DeckSpec(title="述标", slides=[
        {"id": "s1", "title": "报价构成", "kind": "content", "layout": "chart",
         "chart": {"type": "pie", "categories": [f"项{i}" for i in range(9)],
                   "series": [{"name": "金额", "values": [1.0] * 9}]}},
    ])
    prs = Presentation(io.BytesIO(render_pptx(deck)))
    assert any(sh.has_chart for sh in prs.slides[0].shapes)


def test_chart_page_renders_every_bullet_not_just_the_first():
    """提示词允许图表页带 1-2 句结论，编辑器也保存全部要点——只画 bullets[0] 等于用户写的
    第二句在 PPT 里凭空消失且毫无提示（评审）。"""
    deck = DeckSpec(title="述标", slides=[
        {"id": "s1", "title": "团队", "kind": "content", "layout": "chart",
         "bullets": ["第一句结论", "第二句结论"],
         "chart": {"type": "pie", "categories": ["高级"], "series": [{"name": "人数", "values": [3]}]}},
    ])
    prs = Presentation(io.BytesIO(render_pptx(deck)))
    texts = " ".join(_all_texts(prs))
    assert "第一句结论" in texts and "第二句结论" in texts


def test_new_layouts_fit_inside_a_four_by_three_master():
    """企业母版路径沿用客户自己的页面尺寸。新版式若按死写的 16:9 常量绘制，4:3 母版（10in 宽）
    下图表会被推出页面 2.63in（约四分之一被裁），页码则每页都跑到页面外（评审实测复现）。"""
    master = _tiny_master(width_in=10.0, height_in=7.5)
    deck = DeckSpec(title="述标", slides=[
        {"id": "sec", "title": "技术方案", "kind": "section", "bullets": ["过渡语"]},
        {"id": "s1", "title": "业绩", "kind": "content", "layout": "chart",
         "chart": {"type": "column", "categories": ["A", "B"], "series": [{"name": "n", "values": [1, 2]}]}},
        {"id": "s2", "title": "对比", "kind": "content", "layout": "comparison",
         "bullets": ["要点"], "stats": [{"value": "72 小时", "label": "提前完成"}]},
    ])
    prs = Presentation(io.BytesIO(render_pptx(deck, master_bytes=master)))
    assert prs.slide_width == Inches(10.0)
    overflow = [
        sh.left + sh.width
        for sl in prs.slides for sh in sl.shapes
        if sh.left is not None and sh.width is not None and sh.left + sh.width > prs.slide_width + Emu(9144)
    ]
    assert overflow == [], f"有形状超出母版页宽: {overflow}"


def test_chart_renders_even_when_the_model_omitted_the_layout_key():
    """端到端守住那次事故：deck 里带 chart 数据但 layout 缺省成 bullets 时，导出必须仍然渲出
    真实图表对象，而不是把数据丢掉只画要点（用户实际拿到的就是后者）。
    存量 deck 重新导出也走这条路径——库里已有的坏形状导出即自愈，无需重新生成。"""
    deck = DeckSpec(title="述标", slides=[
        {"id": "s1", "title": "项目实施进度与团队配置", "kind": "content",
         "bullets": ["到货验收后 15 日内完成安装调试"],
         "chart": {"type": "column", "categories": ["安装", "调试", "培训"],
                   "series": [{"name": "工期(天)", "values": [5, 7, 3]}]}},
    ])
    assert deck.slides[0].layout == "chart"          # 数据即意图
    prs = Presentation(io.BytesIO(render_pptx(deck)))
    charts = [sh for sh in prs.slides[0].shapes if sh.has_chart]
    assert len(charts) == 1, "图表数据在，导出就必须有真实图表对象"
    assert list(charts[0].chart.plots[0].categories) == ["安装", "调试", "培训"]


def _shapes_of(prs, idx=0):
    return list(prs.slides[idx].shapes)


def test_scoring_chip_is_wide_enough_for_chinese_and_never_bleeds_left():
    """用户实测：每一页左下角的评分点角标渲成「分点｜…」——「评」字被切在画面外。
    _chip_width 按 0.11in/字符估宽，而 12pt 中文字符实际约 0.167in：30 字标签估 3.3in、
    实需 ~4.7in，形状文字默认居中且 word_wrap=False，多出的 1.4in 往两边各溢出 0.7in，
    左边正好顶着页边距。宽度要按字符类型估，且文字左对齐——溢出只能往右，不能吃掉首字。"""
    from pptx.enum.text import PP_ALIGN
    from agent.agents.bidding_agent.render.pptx import _chip_width
    from pptx.util import Inches, Pt

    scoring = "★交货时间/地点/方式、★质保期、★知识产权与保密要求"
    text = f"评分点｜{scoring}"
    cjk = sum(1 for c in text if ord(c) > 0x2E80)
    # 12pt 中文 ≈ 12pt 宽 = 1/6 in；估宽必须覆盖真实文字宽度，否则必然溢出
    assert _chip_width(text) >= Inches(cjk / 6.0), "中文角标估宽不足，文字会溢出框外"

    deck = DeckSpec(title="述标", slides=[
        {"id": "s1", "title": "核心需求", "kind": "content", "bullets": ["要点"], "scoring": scoring},
    ])
    prs = Presentation(io.BytesIO(render_pptx(deck)))
    chips = [sh for sh in _shapes_of(prs)
             if sh.has_text_frame and sh.text_frame.text.startswith("评分点")]
    assert chips, "评分点角标没渲出来"
    chip = chips[0]
    assert chip.left >= 0
    assert chip.text_frame.paragraphs[0].alignment == PP_ALIGN.LEFT, "居中会让溢出吃掉左侧首字"


def test_bullets_fill_the_slide_instead_of_hugging_the_top():
    """用户反馈「太素」的一半原因：4 条要点只占顶部约 25%，下面 60% 是纯白。
    要点区固定 top=1.4in、顶对齐，内容少时下方全空。改成卡片等距铺开并整体垂直居中后，
    最后一张卡片必须越过页面中线，页面才不会头重脚轻。"""
    deck = DeckSpec(title="述标", slides=[
        {"id": "s1", "title": "核心需求理解", "kind": "content", "scoring": "技术方案 15 分",
         "bullets": ["承诺 30 日内交付", "3 年免费质保", "原厂全新正品", "全部★条款无偏离"]},
    ])
    prs = Presentation(io.BytesIO(render_pptx(deck)))
    mid = prs.slide_height / 2
    bodies = [sh for sh in _shapes_of(prs)
              if sh.has_text_frame and "质保" in sh.text_frame.text]
    assert bodies, "要点没渲出来"
    lowest = max(sh.top + sh.height for sh in bodies)
    assert lowest > mid, "要点全挤在上半页，下面一大片空白"


def _deck_for_style():
    return DeckSpec(title="述标", slides=[
        {"id": "c", "title": "封面", "kind": "cover", "bullets": []},
        {"id": "sec", "title": "第一部分 技术方案", "kind": "section", "bullets": []},
        {"id": "s1", "title": "核心需求", "kind": "content", "scoring": "技术方案 20 分",
         "bullets": ["要点一", "要点二", "要点三"]},
        {"id": "e", "title": "感谢聆听", "kind": "end", "bullets": []},
    ])


def _layout_signature(prs: Presentation, idx: int) -> tuple:
    """某一页的版式指纹：所有图形的尺寸集合。只换配色不换版式的话三套会撞成同一个指纹。"""
    return tuple(sorted((sh.width, sh.height) for sh in prs.slides[idx].shapes
                        if sh.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE))


def test_three_templates_differ_in_layout_not_only_colour():
    """模板此前只差三个配色、版式完全一样，用户选来选去每页长得一模一样，等于没得选。
    这条锁住「换模板必须换版式」：封面（满幅/分栏/横幅）与正文页（overline/numeral/corner
    标题 + numbered/hairline/elevated 卡片）在三套之间都不能撞。"""
    covers, contents = {}, {}
    for tpl in ("blue", "gov", "tech"):
        prs = Presentation(io.BytesIO(render_pptx(_deck_for_style(), template=tpl)))
        covers[tpl] = _layout_signature(prs, 0)
        contents[tpl] = _layout_signature(prs, 2)
    assert len(set(covers.values())) == 3, f"三套封面版式没有区分：{covers}"
    assert len(set(contents.values())) == 3, f"三套正文版式没有区分：{contents}"


def test_unknown_layout_switches_fall_back_to_the_default_painters(monkeypatch):
    """结构开关是枚举：token 里写了渲染层不认识的取值（改名/手写脏数据/未来新模板漏登记画法），
    必须静默回退到默认那套画法，而不是 KeyError 让整份述标导不出来。"""
    from agent.agents.bidding_agent.render import styles
    bogus = dict(_TEMPLATE_TOKENS["blue"], header="???", card="???", cover="???")
    monkeypatch.setitem(styles.TEMPLATE_TOKENS, "bogus", bogus)
    prs = Presentation(io.BytesIO(render_pptx(_deck_for_style(), template="bogus")))
    ref = Presentation(io.BytesIO(render_pptx(_deck_for_style(), template="blue")))
    assert len(prs.slides) == len(ref.slides)
    for idx in range(len(ref.slides)):
        assert _layout_signature(prs, idx) == _layout_signature(ref, idx), f"第 {idx} 页没有回退到默认版式"


def _rel_luminance(rgb: RGBColor) -> float:
    def channel(v: int) -> float:
        c = v / 255
        return c / 12.92 if c <= 0.03928 else ((c + 0.055) / 1.055) ** 2.4
    r, g, b = rgb
    return 0.2126 * channel(r) + 0.7152 * channel(g) + 0.0722 * channel(b)


def _contrast(a: RGBColor, b: RGBColor) -> float:
    la, lb = _rel_luminance(a), _rel_luminance(b)
    return (max(la, lb) + 0.05) / (min(la, lb) + 0.05)


def test_every_template_keeps_its_text_readable_on_its_own_ground():
    """三套模板的每一对「字色 × 它实际压着的底色」都要过 4.5:1（正文/弱化字/主色块反白/
    卡片与角标）。投影仪比屏幕更不宽容，浅底那套颜色搬到深底上就是一团糊——深色模板的
    on_primary 不是白色正是为了这条。"""
    from agent.agents.bidding_agent.render.styles import on_primary
    for name, t in _TEMPLATE_TOKENS.items():
        pairs = {
            "正文/页底": (t["text"], t["bg"]),
            "弱化字/页底": (t["muted"], t["bg"]),
            "主色块反白": (on_primary(t), t["primary"]),
            "强调色/卡片底": (t["accent"], t["tint"]),
            "正文/卡片底": (t["text"], t["tint"]),
        }
        for label, (fg, bg) in pairs.items():
            ratio = _contrast(fg, bg)
            assert ratio >= 4.5, f"{name} 的「{label}」对比度只有 {ratio:.2f}:1，投影上会糊"


def test_dark_template_paints_a_background_and_light_text():
    """深色模板必须真的铺底色并把正文改成浅色——只改强调色的话，深色主题的字仍是深灰，
    在白底上看不出区别，在深底上则糊成一片。"""
    prs = Presentation(io.BytesIO(render_pptx(_deck_for_style(), template="tech")))
    content = prs.slides[2]
    slide_w, slide_h = prs.slide_width, prs.slide_height
    full = [sh for sh in content.shapes
            if sh.width == slide_w and sh.height == slide_h and sh.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE]
    assert full, "深色模板没有铺满整页的底色块"
    assert full[0].fill.fore_color.rgb == RGBColor(15, 23, 42)
    body = [sh for sh in content.shapes if sh.has_text_frame and sh.text_frame.text == "要点一"]
    assert body, "要点没渲出来"
    colour = body[0].text_frame.paragraphs[0].runs[0].font.color.rgb
    assert colour == RGBColor(226, 232, 240), f"深底上的正文色仍是浅底那套：{colour}"


def test_light_templates_do_not_paint_a_background():
    """浅色模板不铺整页底色：空白版式本来就是白底，多画一层只会让文件变大、还可能盖住母版元素。
    （封面是例外：满幅/横幅封面本来就要一块自己的底，所以只看正文页。）"""
    for tpl in ("blue", "gov"):
        prs = Presentation(io.BytesIO(render_pptx(_deck_for_style(), template=tpl)))
        content = prs.slides[2]
        full = [sh for sh in content.shapes
                if sh.width == prs.slide_width and sh.height == prs.slide_height]
        assert not full, f"{tpl} 不该铺整页底色"


def _full_deck():
    """一份把所有版式都用上的 deck：封面/分隔页/要点页/饼图页/多系列柱图页/对比页/结束页。"""
    return DeckSpec(title="某市政务云平台运维服务项目", slides=[
        {"id": "c", "kind": "cover", "title": "某市政务云平台运维服务项目",
         "bullets": ["投标人：某科技股份有限公司", "述标时长 15 分钟"]},
        {"id": "sec", "kind": "section", "title": "项目理解与总体方案", "bullets": ["逐条对应评分办法"]},
        {"id": "b", "kind": "content", "title": "核心需求理解", "scoring": "★服务范围 20 分",
         "bullets": ["7×24 驻场值守", "一级故障 2 小时恢复", "全部★条款无偏离"], "notes": "口播稿"},
        {"id": "p", "kind": "content", "title": "团队构成", "layout": "chart", "scoring": "团队 15 分",
         "bullets": ["中级及以上职称占比 60%"],
         "chart": {"type": "pie", "categories": ["高级", "中级", "初级"],
                   "series": [{"name": "人数", "values": [3, 6, 4]}]}},
        {"id": "col", "kind": "content", "title": "三年投入", "layout": "chart",
         "chart": {"type": "column", "categories": ["2024", "2025", "2026"],
                   "series": [{"name": "人月", "values": [96, 132, 156]},
                              {"name": "巡检", "values": [48, 60, 72]}]}},
        {"id": "cmp", "kind": "content", "title": "承诺优于要求", "layout": "comparison",
         "scoring": "★服务承诺 15 分", "bullets": ["恢复时限提前 50%", "质保 3 年"],
         "stats": [{"value": "2 小时", "label": "一级故障恢复"}, {"value": "0 起", "label": "质量投诉"}]},
        {"id": "e", "kind": "end", "title": "感谢聆听"},
    ])


def test_every_template_renders_every_layout_and_reopens():
    """三套 × 全版式（图表/对比/要点/分隔/封面/结束）都要能渲出来、能被重新打开、每页有形状。
    这是模板改版的兜底回归：换了骨架不能有哪一套在某个版式上画不出东西或直接抛异常。"""
    for tpl in ("blue", "tech", "gov"):
        data = render_pptx(_full_deck(), template=tpl)
        assert data[:2] == b"PK"
        prs = Presentation(io.BytesIO(data))
        assert len(prs.slides) == 7
        for i, slide in enumerate(prs.slides):
            assert len(slide.shapes) > 0, f"{tpl} 第 {i} 页是空页"
        assert sum(1 for sl in prs.slides for sh in sl.shapes if sh.has_chart) == 2
        texts = " ".join(_all_texts(prs))
        assert "核心需求理解" in texts and "承诺优于要求" in texts


def _master_without_placeholders(width_in: float = 10.0, height_in: float = 7.5, *,
                                 every_layout: bool = False) -> bytes:
    """把首个版式的占位符全部摘掉的迷你母版：客户母版里这种「纯图形版式」很常见，
    此时封面/结束页会退回我们自己的画法——新封面骨架正是在这条路径上才会被画出来。
    every_layout=True 连正文版式的占位符一起摘掉，正文页的要点也会退回我们自己的密度版式。"""
    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(width_in), Inches(height_in)
    layouts = list(prs.slide_layouts) if every_layout else [prs.slide_layouts[0]]
    for layout in layouts:
        for ph in list(layout.placeholders):
            ph._element.getparent().remove(ph._element)
    out = io.BytesIO()
    prs.save(out)
    return out.getvalue()


def test_covers_fall_back_to_our_own_painting_on_a_placeholderless_master():
    """母版没有标题占位符时，封面/结束页整页退回我们自己的画法——三套的新封面骨架
    （满幅/分栏/横幅）都得在客户自己的页面尺寸里画完，一寸都不能出界。"""
    master = _master_without_placeholders(width_in=10.0, height_in=7.5)
    for tpl in ("blue", "tech", "gov"):
        prs = Presentation(io.BytesIO(render_pptx(_full_deck(), template=tpl, master_bytes=master)))
        cover = prs.slides[0]
        assert any(sh.has_text_frame and "某市政务云平台运维服务项目" in sh.text_frame.text
                   for sh in cover.shapes), f"{tpl} 封面没画出来"
        slop = Emu(9144)
        over = [(sh.left, sh.width) for sh in cover.shapes
                if sh.left + sh.width > prs.slide_width + slop or sh.left < -slop]
        assert over == [], f"{tpl} 的封面在窄母版上画出界：{over}"


def test_every_template_fits_inside_a_narrow_enterprise_master():
    """企业母版路径 × 三套模板：新骨架（角标/分栏封面/横幅封面）一律按真实页宽页高定位，
    在 4:3 母版（10in 宽）上不许有任何形状越过页边——死写 16:9 常量就会每页被裁掉一截。"""
    master = _tiny_master(width_in=10.0, height_in=7.5)
    for tpl in ("blue", "tech", "gov"):
        prs = Presentation(io.BytesIO(render_pptx(_full_deck(), template=tpl, master_bytes=master)))
        assert prs.slide_width == Inches(10.0)
        assert len(prs.slides) == 7
        slop = Emu(9144)      # 0.01in：形状边框宽度级别的容差
        over = [(i, sh.shape_type, sh.left + sh.width, sh.top + sh.height)
                for i, sl in enumerate(prs.slides) for sh in sl.shapes
                if sh.left is not None and sh.width is not None
                and (sh.left + sh.width > prs.slide_width + slop
                     or sh.top + sh.height > prs.slide_height + slop
                     or sh.left < -slop or sh.top < -slop)]
        assert over == [], f"{tpl} 在 4:3 母版上画出界：{over}"


# ---- 视觉精致度第二轮：图表页双栏 / 要点密度 / 字号层级 / 圆角常量 / 饼图色阶 / 封面锚点 ----

def _content_geometry(prs):
    """(正文左边距, 正文可用宽)：与渲染层同一套几何。"""
    from agent.agents.bidding_agent.render.pptx_blocks import MARGIN
    return MARGIN, prs.slide_width - 2 * MARGIN


def _chart_deck(bullets: list[str]) -> DeckSpec:
    return DeckSpec(title="述标", slides=[
        {"id": "s1", "title": "团队构成", "kind": "content", "layout": "chart", "bullets": bullets,
         "chart": {"type": "pie", "categories": ["高级", "中级", "初级"],
                   "series": [{"name": "人数", "values": [3, 6, 4]}]}},
    ])


def test_chart_page_pairs_the_chart_with_a_conclusions_rail():
    """图表页有自己的双栏版式：图表占主区，右边一条「关键结论」栏。
    此前图表卡在正文带里、结论压成一行小字贴在图下沿，左右各空出一大片——饼图尤其明显
    （直径被页高卡住，横向撑不到 12in），用户实评「图表页左右偏空」。"""
    prs = Presentation(io.BytesIO(render_pptx(_chart_deck(["中级及以上占比 60%", "持证率 100%"]))))
    slide = prs.slides[0]
    left, content_w = _content_geometry(prs)
    chart = next(sh for sh in slide.shapes if sh.has_chart)
    assert chart.left == left
    assert chart.width <= content_w * 0.7, "图表仍霸着整幅宽，双栏没生效"
    assert chart.height > prs.slide_height * 0.55, "图表没吃满正文带的高度"
    rail = next(sh for sh in slide.shapes if sh.has_text_frame and sh.text_frame.text == "关键结论")
    assert rail.left > chart.left + chart.width, "结论栏没落在图表右边"
    panel = next(sh for sh in slide.shapes
                 if sh.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE
                 and sh.auto_shape_type == MSO_SHAPE.ROUNDED_RECTANGLE
                 and sh.left > chart.left + chart.width)
    # 两栏合起来正好铺满正文可用宽：右边不再剩一条没人用的白
    assert abs((panel.left + panel.width) - (left + content_w)) <= Emu(9144)
    texts = " ".join(_all_texts(prs))
    assert "中级及以上占比 60%" in texts and "持证率 100%" in texts


def test_chart_page_without_conclusions_lets_the_chart_span_the_whole_width():
    """没有结论可放就不画侧栏、也不编内容，让图表自己占满整幅——空栏比空白更难看。"""
    prs = Presentation(io.BytesIO(render_pptx(_chart_deck([]))))
    _, content_w = _content_geometry(prs)
    chart = next(sh for sh in prs.slides[0].shapes if sh.has_chart)
    assert chart.width == content_w
    assert "关键结论" not in " ".join(_all_texts(prs)), "没有结论却画出了空侧栏"


def _bullet_layout(n: int, *, template: str = "blue", layout: str = "bullets") -> tuple:
    """n 条要点渲出来的 (列数, 正文字号, 卡片高度)。"""
    items = [f"要点{i + 1}" for i in range(n)]
    slide = {"id": "s1", "title": "要点页", "kind": "content", "bullets": items, "layout": layout}
    if layout == "comparison":
        slide["stats"] = [{"value": "2 小时", "label": "恢复"}]
    prs = Presentation(io.BytesIO(render_pptx(DeckSpec(title="述标", slides=[slide]),
                                              template=template)))
    boxes = [sh for sh in prs.slides[0].shapes
             if sh.has_text_frame and sh.text_frame.text in items]
    assert len(boxes) == n, "要点没有逐条渲出来"
    return (len({sh.left for sh in boxes}),
            boxes[0].text_frame.paragraphs[0].runs[0].font.size.pt,
            max(sh.height for sh in boxes))


def test_bullet_pages_switch_layout_with_the_item_count():
    """2 条和 6 条不能长一个样（用户实评：2 条空得慌、6 条挤）——
    ≤3 条大卡纵向、4-6 条双栏、>6 条紧凑行，字号与卡高跟着退让。"""
    two, five, eight = _bullet_layout(2), _bullet_layout(5), _bullet_layout(8)
    assert (two[0], five[0], eight[0]) == (1, 2, 1), f"列数没随条数变：{two} {five} {eight}"
    assert two[1] > five[1] > eight[1], f"字号没随密度退让：{two} {five} {eight}"
    assert two[2] > five[2] > eight[2], f"卡高没随密度收紧：{two} {five} {eight}"


def test_a_narrow_column_never_splits_into_two():
    """对比页左栏只有 56% 宽，4-6 条再切两半每张卡放不下一句话——窄栏一律退回单列。"""
    assert _bullet_layout(5, layout="comparison")[0] == 1


def test_rendered_type_scale_keeps_title_body_and_note_apart():
    """标题/正文/注释三档必须靠字号 × 字重 × 色阶同时拉开：三档挨在一起就是
    「排得整齐但不惊艳」。断言读渲染产物而不是常量表——常量改了画法没跟上，这条要红。"""
    deck = DeckSpec(title="某项目", slides=[
        {"id": "s1", "title": "核心需求理解", "kind": "content", "scoring": "服务 20 分",
         "bullets": ["要点一", "要点二", "要点三", "要点四"]},
    ])
    prs = Presentation(io.BytesIO(render_pptx(deck)))

    def _run(pred):
        sh = next(s for s in prs.slides[0].shapes if s.has_text_frame and pred(s.text_frame.text))
        return sh.text_frame.paragraphs[0].runs[0]

    title = _run(lambda t: t == "核心需求理解")
    body = _run(lambda t: t == "要点一")
    note = _run(lambda t: t.startswith("评分点"))
    caption = _run(lambda t: "/" in t and not t.startswith("评分点"))
    assert title.font.size.pt >= body.font.size.pt * 1.8, "标题不够大，压不住正文"
    assert title.font.bold is True and not body.font.bold, "字重没有拉开"
    assert note.font.size.pt <= body.font.size.pt * 0.85, "注释没有明显退到正文后面"
    assert caption.font.size.pt <= note.font.size.pt


def _rounded(prs):
    return [sh for sl in prs.slides for sh in sl.shapes
            if sh.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE
            and sh.auto_shape_type == MSO_SHAPE.ROUNDED_RECTANGLE]


def test_every_rounded_card_shares_one_corner_radius():
    """圆角要统一到一个**绝对**半径：默认圆角是短边的 16.7%，于是 0.5in 高的角标
    和 3in 高的数字卡圆角差着六倍——同一页好几种圆角，一眼就是默认值堆出来的。"""
    for tpl in ("blue", "tech", "gov"):
        prs = Presentation(io.BytesIO(render_pptx(_full_deck(), template=tpl)))
        shapes = _rounded(prs)
        assert len(shapes) >= 5, f"{tpl} 没渲出圆角卡片"
        for sh in shapes:
            radius = sh.adjustments[0] * min(sh.width, sh.height)
            assert abs(radius - CARD["radius"]) <= Emu(9144), \
                f"{tpl} 有圆角半径 {radius} 偏离统一常量 {CARD['radius']}"


def _lab(rgb):
    def gamma(c: int) -> float:
        v = c / 255
        return v / 12.92 if v <= 0.04045 else ((v + 0.055) / 1.055) ** 2.4

    def f(t: float) -> float:
        return t ** (1 / 3) if t > 0.008856 else 7.787 * t + 16 / 116
    r, g, b = (gamma(c) for c in rgb)
    x = f((0.4124 * r + 0.3576 * g + 0.1805 * b) / 0.95047)
    y = f(0.2126 * r + 0.7152 * g + 0.0722 * b)
    z = f((0.0193 * r + 0.1192 * g + 0.9505 * b) / 1.08883)
    return (116 * y - 16, 500 * (x - y), 200 * (y - z))


def _delta_e(a, b) -> float:
    return sum((p - q) ** 2 for p, q in zip(_lab(a), _lab(b))) ** 0.5


def test_pie_palette_covers_eight_slices_distinctly_and_readably():
    """>4 类的饼图色阶：8 类不许重复取色（上一轮的短板——封顶之后相邻两块同色），
    每一块还要压得住那唯一一种数值标签字色（渲染器不认逐扇区字色，只能整图一种）。"""
    from agent.agents.bidding_agent.render.styles import chart_palette, on_primary
    for name, tokens in _TEMPLATE_TOKENS.items():
        label = on_primary(tokens)
        colors = chart_palette(tokens, 8, label_rgb=label)
        assert len({str(c) for c in colors}) == 8, f"{name} 的 8 类色阶有重复色：{colors}"
        for c in colors:
            assert _contrast(c, label) >= 4.5, f"{name} 的扇区 {c} 压不住数值标签"
        gaps = [_delta_e(colors[i], colors[j]) for i in range(8) for j in range(i + 1, 8)]
        assert min(gaps) >= 9.0, f"{name} 有两块几乎同色（ΔE 仅 {min(gaps):.1f}）"


def test_eight_slice_pie_really_paints_eight_different_colours():
    """色阶算得对还不够——真渲出来的 8 个扇区必须是 8 种颜色。"""
    deck = DeckSpec(title="述标", slides=[
        {"id": "s1", "title": "报价构成", "kind": "content", "layout": "chart",
         "chart": {"type": "pie", "categories": [f"项{i}" for i in range(8)],
                   "series": [{"name": "金额", "values": [float(i + 1) for i in range(8)]}]}},
    ])
    prs = Presentation(io.BytesIO(render_pptx(deck, template="blue")))
    chart = next(sh for sh in prs.slides[0].shapes if sh.has_chart).chart
    fills = {str(p.format.fill.fore_color.rgb) for p in chart.plots[0].series[0].points}
    assert len(fills) == 8, f"饼图实际画出来只有 {len(fills)} 种扇区色"


def _hairlines(prs, *, vertical: bool) -> list:
    """封面上的"细线"形状：一个方向很细、另一个方向够长的矩形。"""
    sw, sh = prs.slide_width, prs.slide_height
    out = []
    for s in prs.slides[0].shapes:
        if s.shape_type != MSO_SHAPE_TYPE.AUTO_SHAPE:
            continue
        thin, long_, span = (s.width, s.height, sh) if vertical else (s.height, s.width, sw)
        if thin <= Emu(63500) and long_ >= span * 0.2:
            out.append(s)
    return out


def test_each_cover_carries_its_own_visual_anchor():
    """三套封面各有一处克制的几何锚点（栏位竖线节奏 / 网格 + 裁切角 / 内缩烫金边框）：
    只有一块底色加一行标题的封面立不住，用户实评「偏素」。"""
    blue = Presentation(io.BytesIO(render_pptx(_deck(), template="blue")))
    rhythm = [s for s in _hairlines(blue, vertical=True) if s.left > blue.slide_width * 0.5]
    assert len(rhythm) >= 4, "商务提案封面右半页没有栏位竖线节奏"

    tech = Presentation(io.BytesIO(render_pptx(_deck(), template="tech")))
    seam = tech.slide_width * 0.45
    dots = [s for s in tech.slides[0].shapes
            if s.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE and s.left >= seam
            and s.width <= Emu(63500) and s.height <= Emu(63500)]
    assert len(dots) >= 12, "技术方案封面右栏没有点阵锚点"

    gov = Presentation(io.BytesIO(render_pptx(_deck(), template="gov")))
    frame_h = [s for s in _hairlines(gov, vertical=False)
               if s.left > 0 and s.left + s.width < gov.slide_width]
    frame_v = [s for s in _hairlines(gov, vertical=True)
               if s.top > 0 and s.top + s.height < gov.slide_height]
    assert len(frame_h) >= 2 and len(frame_v) >= 2, "党政庄重封面没有内缩一圈的烫金边框"


def test_text_boxes_align_optically_with_the_shapes_beside_them():
    """光学边距：文本框默认左右各留 0.1in 内边距，于是「放在 x 处的文字」实际画在 x+0.1in——
    同一条左边线上的色块和标题就差了 0.1in，整页看着像没对齐。我们的文本框四边内边距一律清零，
    需要留白就把留白算进几何里；页码同理，右沿要落在版心右边距上而不是随手贴近页边。"""
    prs = Presentation(io.BytesIO(render_pptx(_deck(), template="blue")))
    content = prs.slides[1]
    title = next(sh for sh in content.shapes
                 if sh.has_text_frame and sh.text_frame.text == "运维体系")
    rule = next(sh for sh in content.shapes
                if sh.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE
                and _solid_rgb(sh) == _TEMPLATE_TOKENS["blue"]["accent"]
                and sh.top < prs.slide_height * 0.2)
    assert title.left == rule.left, "标题框与它上方的强调短线不在同一条左边线上"
    assert title.text_frame.margin_left == 0, "文本框还留着默认内边距，文字比色块右移 0.1in"
    left, content_w = _content_geometry(prs)
    page_no = next(sh for sh in content.shapes
                   if sh.has_text_frame and sh.text_frame.text == "1 / 2")
    assert abs((page_no.left + page_no.width) - (left + content_w)) <= Emu(9144), \
        "页码右沿没有落在版心右边距上"


def test_dense_layouts_fit_inside_a_narrow_enterprise_master():
    """双栏要点 / 紧凑行 / 图表页结论侧栏这三种新版式，在 4:3 母版（10in 宽）里也要画得下。
    双栏槽宽与侧栏比例都是按**正文可用宽**算的；任何一处退回 16:9 常量，窄母版上就会出界。
    母版特意挑「占位符被摘光」的那种：有正文占位符时要点走母版自己的占位符，
    我们的密度版式根本不会被调用，测了个寂寞。"""
    master = _master_without_placeholders(width_in=10.0, height_in=7.5, every_layout=True)
    deck = DeckSpec(title="述标", slides=[
        {"id": "d1", "kind": "content", "title": "五项措施",
         "bullets": [f"措施{i + 1}" for i in range(5)]},
        {"id": "d2", "kind": "content", "title": "八项投入",
         "bullets": [f"投入{i + 1}" for i in range(8)]},
        {"id": "d3", "kind": "content", "title": "报价构成", "layout": "chart",
         "bullets": ["人力成本占比 46%"],
         "chart": {"type": "pie", "categories": [f"项{i}" for i in range(8)],
                   "series": [{"name": "金额", "values": [float(i + 1) for i in range(8)]}]}},
    ])
    for tpl in ("blue", "tech", "gov"):
        prs = Presentation(io.BytesIO(render_pptx(deck, template=tpl, master_bytes=master)))
        slop = Emu(9144)
        over = [(i, sh.left, sh.width) for i, sl in enumerate(prs.slides) for sh in sl.shapes
                if sh.left is not None and sh.width is not None
                and (sh.left + sh.width > prs.slide_width + slop or sh.left < -slop)]
        assert over == [], f"{tpl} 的密集版式在 4:3 母版上画出界：{over}"
